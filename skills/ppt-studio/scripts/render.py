#!/usr/bin/env python3
"""render.py — dual-render PPT engine.

One plan.json drives two outputs that share identical layout logic:

  render.py --plan plan.json --preview preview.html
      Generate a standalone HTML preview (all slides stacked).

  render.py --plan plan.json --rasterize qa/
      Screenshot each slide from its HTML into qa/slide-N.jpg (for QA).

  render.py --plan plan.json --export output.pptx
      Generate the final editable .pptx.

Architecture: see references/generate/render-architecture.md
"""

import argparse
import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from ppt_render.core import (
    build_preview_html, single_slide_html, render_slide_to_pptx,
    CANVAS_W, CANVAS_H, PREVIEW_SCALE,
)
from ppt_render.templates import TEMPLATES

DEFAULT_THEME = {
    "primary": "2A3D52", "secondary": "8B97A5", "accent": "FFFFFF",
    "text_dark": "2A3D52", "text_body": "4F5868", "bg_light": "F6F7F9",
    "fontHeading": "Microsoft YaHei", "fontBody": "Microsoft YaHei",
}


def load_plan(path):
    with open(path, "r", encoding="utf-8") as f:
        plan = json.load(f)
    theme = dict(DEFAULT_THEME)
    theme.update(plan.get("theme", {}))
    slides = plan.get("slides", [])
    if not slides:
        raise ValueError("plan.json has no slides")
    return theme, slides


def layout_all(theme, slides):
    """Run every slide's layout() function. Returns (elements_list, labels)."""
    all_elements = []
    labels = []
    for idx, slide in enumerate(slides):
        tname = slide.get("template")
        if tname not in TEMPLATES:
            raise ValueError(
                f"slide {idx + 1}: unknown template '{tname}'. "
                f"Valid: {', '.join(sorted(TEMPLATES))}")
        # ensure page number defaults to position
        if "page" not in slide:
            slide["page"] = idx + 1
        try:
            els = TEMPLATES[tname](slide, theme)
        except Exception as e:
            raise ValueError(f"slide {idx + 1} ({tname}): {e}") from e
        all_elements.append(els)
        labels.append(f"Slide {idx + 1}: {tname}")
    return all_elements, labels


# ── --preview ──────────────────────────────────────────────

def cmd_preview(theme, slides, out_path):
    all_elements, labels = layout_all(theme, slides)
    html = build_preview_html(all_elements, labels)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"Preview written: {out_path} ({len(slides)} slides)")


# ── --rasterize ────────────────────────────────────────────

def cmd_rasterize(theme, slides, out_dir):
    from playwright.sync_api import sync_playwright

    all_elements, _ = layout_all(theme, slides)
    os.makedirs(out_dir, exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": CANVAS_W, "height": CANVAS_H},
            device_scale_factor=PREVIEW_SCALE,
        )
        for i, els in enumerate(all_elements):
            html = single_slide_html(els)
            page.set_content(html, wait_until="networkidle")
            out = os.path.join(out_dir, f"slide-{i + 1}.jpg")
            page.screenshot(path=out, type="jpeg", quality=92,
                            clip={"x": 0, "y": 0,
                                  "width": CANVAS_W, "height": CANVAS_H})
            print(f"  rasterized slide {i + 1} -> {out}")
        browser.close()
    print(f"Rasterized {len(all_elements)} slides into {out_dir}")


# ── --export ───────────────────────────────────────────────

def cmd_export(theme, slides, out_path):
    from pptx import Presentation
    from pptx.util import Inches

    all_elements, _ = layout_all(theme, slides)

    prs = Presentation()
    prs.slide_width = Inches(CANVAS_W / 96)
    prs.slide_height = Inches(CANVAS_H / 96)
    blank = prs.slide_layouts[6]  # blank layout

    for els in all_elements:
        slide = prs.slides.add_slide(blank)
        render_slide_to_pptx(slide, els)

    prs.save(out_path)
    print(f"PPTX written: {out_path} ({len(slides)} slides)")


# ── CLI ────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description="Dual-render PPT engine")
    ap.add_argument("--plan", required=True, help="path to plan.json")
    ap.add_argument("--preview", help="output path for HTML preview")
    ap.add_argument("--rasterize", help="output directory for slide JPGs")
    ap.add_argument("--export", help="output path for .pptx")
    args = ap.parse_args()

    if not (args.preview or args.rasterize or args.export):
        ap.error("specify at least one of --preview / --rasterize / --export")

    try:
        theme, slides = load_plan(args.plan)
        if args.preview:
            cmd_preview(theme, slides, args.preview)
        if args.rasterize:
            cmd_rasterize(theme, slides, args.rasterize)
        if args.export:
            cmd_export(theme, slides, args.export)
    except (ValueError, FileNotFoundError, json.JSONDecodeError) as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
