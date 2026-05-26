#!/usr/bin/env python3
"""
pptx_insert.py — splice new pages from one PPTX into another.

Used by Path G' in the ppt-studio skill. Typical flow:
  1. Path G generates new_pages.pptx (1-N new slides)
  2. This script merges them into the user's original deck

Usage:
    python pptx_insert.py \\
        --original <original.pptx> \\
        --new-pages <new_pages.pptx> \\
        --mode {insert|replace|append} \\
        [--at <position>] \\
        --output <output.pptx>

Modes:
    insert     Insert new pages at --at position (1-based). Original pages shift down.
    replace    Replace the slide at --at with the new pages.
    append     Append new pages at end. --at is ignored.

Examples:
    # Append new pages to end of original
    python pptx_insert.py --original deck.pptx --new-pages new.pptx \\
        --mode append --output deck_extended.pptx

    # Insert at position 3 (becomes new slide 3, old 3+ shifts down)
    python pptx_insert.py --original deck.pptx --new-pages new.pptx \\
        --mode insert --at 3 --output deck_extended.pptx

    # Replace slide 5 with the new pages
    python pptx_insert.py --original deck.pptx --new-pages new.pptx \\
        --mode replace --at 5 --output deck_revised.pptx
"""

import argparse
import io
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# Shape cloning recipes (from references/insert/insert-recipes.md)
# ============================================================

def copy_run_format(src_run, dst_run):
    """Copy font name/size/color/bold/italic from src run to dst run."""
    sf, df = src_run.font, dst_run.font
    if sf.name is not None:
        df.name = sf.name
    if sf.size is not None:
        df.size = sf.size
    if sf.bold is not None:
        df.bold = sf.bold
    if sf.italic is not None:
        df.italic = sf.italic
    try:
        if sf.color.rgb is not None:
            df.color.rgb = sf.color.rgb
    except (AttributeError, ValueError):
        pass


def clone_picture(source_shape, target_slide, verbose=False):
    """Clone a PICTURE shape across PPTXs, including the image bytes."""
    image = source_shape.image
    if verbose:
        print(f"      clone_picture: {len(image.blob)} bytes, ext={image.ext}",
              file=sys.stderr)
    image_stream = io.BytesIO(image.blob)
    new_pic = target_slide.shapes.add_picture(
        image_stream,
        left=source_shape.left, top=source_shape.top,
        width=source_shape.width, height=source_shape.height,
    )
    return new_pic


def clone_text_shape(source_shape, target_slide, verbose=False):
    """Clone a TEXT_BOX with its paragraphs and per-run formatting."""
    new_tb = target_slide.shapes.add_textbox(
        left=source_shape.left, top=source_shape.top,
        width=source_shape.width, height=source_shape.height,
    )
    new_tf = new_tb.text_frame
    try:
        new_tf.word_wrap = source_shape.text_frame.word_wrap
    except (AttributeError, ValueError):
        pass

    # Remove the default empty paragraph
    for p in list(new_tf.paragraphs):
        p._p.getparent().remove(p._p)

    # Copy paragraphs and runs
    for src_p in source_shape.text_frame.paragraphs:
        new_p = new_tf.add_paragraph()
        if src_p._pPr is not None:
            new_p._p.insert(0, deepcopy(src_p._pPr))
        try:
            new_p.alignment = src_p.alignment
        except (AttributeError, ValueError):
            pass
        for src_r in src_p.runs:
            new_r = new_p.add_run()
            new_r.text = src_r.text
            copy_run_format(src_r, new_r)
    return new_tb


def clone_auto_shape(source_shape, target_slide, verbose=False):
    """Clone an AUTO_SHAPE (rectangle, oval, etc.) preserving fill and geometry."""
    try:
        new_shape = target_slide.shapes.add_shape(
            autoshape_type_id=source_shape.auto_shape_type,
            left=source_shape.left, top=source_shape.top,
            width=source_shape.width, height=source_shape.height,
        )
    except (AttributeError, ValueError) as e:
        if verbose:
            print(f"      clone_auto_shape: fallback to generic XML copy ({e})",
                  file=sys.stderr)
        return clone_generic_xml(source_shape, target_slide, verbose)

    # Copy fill
    try:
        if source_shape.fill.type == 1:  # SOLID
            new_shape.fill.solid()
            new_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
    except (AttributeError, ValueError):
        pass

    # Copy line
    try:
        if source_shape.line.color.rgb is not None:
            new_shape.line.color.rgb = source_shape.line.color.rgb
        else:
            new_shape.line.fill.background()
    except (AttributeError, ValueError):
        pass

    # Copy text if any
    if source_shape.has_text_frame and source_shape.text_frame.text.strip():
        new_tf = new_shape.text_frame
        for p in list(new_tf.paragraphs):
            p._p.getparent().remove(p._p)
        for src_p in source_shape.text_frame.paragraphs:
            new_p = new_tf.add_paragraph()
            if src_p._pPr is not None:
                new_p._p.insert(0, deepcopy(src_p._pPr))
            for src_r in src_p.runs:
                new_r = new_p.add_run()
                new_r.text = src_r.text
                copy_run_format(src_r, new_r)

    return new_shape


def clone_generic_xml(source_shape, target_slide, verbose=False):
    """Last-resort: deepcopy the shape's XML. Only safe when no external part refs."""
    new_el = deepcopy(source_shape._element)
    target_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")


def clone_shape_to_slide(source_shape, target_slide, verbose=False):
    """Dispatch a single shape clone based on its type."""
    st = source_shape.shape_type

    if st == MSO_SHAPE_TYPE.PICTURE:
        clone_picture(source_shape, target_slide, verbose)
    elif st == MSO_SHAPE_TYPE.GROUP:
        # Recurse — children become top-level on target (we don't recreate the group)
        for child in source_shape.shapes:
            clone_shape_to_slide(child, target_slide, verbose)
    elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        clone_auto_shape(source_shape, target_slide, verbose)
    elif source_shape.has_text_frame:
        # TEXT_BOX or other text-bearing non-auto shapes
        clone_text_shape(source_shape, target_slide, verbose)
    else:
        if verbose:
            print(f"      clone fallback: type={st}", file=sys.stderr)
        clone_generic_xml(source_shape, target_slide, verbose)


# ============================================================
# Slide-level operations
# ============================================================

def pick_blank_layout(prs):
    """Pick the blank-est available layout in the target presentation."""
    for layout in prs.slide_layouts:
        if layout.name in ("Blank", "空白"):
            return layout
    return prs.slide_layouts[-1]


def copy_slide_into_deck(source_slide, target_prs, verbose=False):
    """Copy a slide from one PPTX into another (always appended at end first)."""
    blank_layout = pick_blank_layout(target_prs)
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Strip default placeholders
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    # Clone all source shapes onto new_slide
    for source_shape in source_slide.shapes:
        clone_shape_to_slide(source_shape, new_slide, verbose)

    return new_slide


def get_slide_id_ns():
    """Return the OOXML namespace key for slide rId attribute."""
    return "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def move_appended_to_position(prs, num_appended, target_position):
    """After appending N slides, move them as a block to target_position (1-based)."""
    slide_id_list = prs.slides._sldIdLst
    all_ids = list(slide_id_list)
    initial_count = len(all_ids) - num_appended

    if num_appended == 0:
        return

    # Slides we just added (at the tail)
    appended_ids = all_ids[initial_count:]

    # Remove them from end
    for sid in appended_ids:
        slide_id_list.remove(sid)

    # Insert at target position (1-based → 0-based)
    insert_idx = target_position - 1
    for offset, sid in enumerate(appended_ids):
        slide_id_list.insert(insert_idx + offset, sid)


def delete_slide_at(prs, position):
    """Delete the slide at 1-based position."""
    slide_id_list = prs.slides._sldIdLst
    all_ids = list(slide_id_list)
    target_idx = position - 1
    if not (0 <= target_idx < len(all_ids)):
        raise ValueError(
            f"position {position} out of range (deck has {len(all_ids)} slides)"
        )
    target_sid = all_ids[target_idx]
    rId = target_sid.attrib[get_slide_id_ns()]
    prs.part.drop_rel(rId)
    slide_id_list.remove(target_sid)


# ============================================================
# Top-level mode handlers
# ============================================================

def do_append(original_path, new_pages_path, output_path, verbose=False):
    original = Presentation(original_path)
    new_deck = Presentation(new_pages_path)

    for idx, new_slide in enumerate(new_deck.slides, 1):
        if verbose:
            print(f"  appending slide {idx}/{len(new_deck.slides)}", file=sys.stderr)
        copy_slide_into_deck(new_slide, original, verbose)

    original.save(output_path)
    return len(new_deck.slides)


def do_insert(original_path, new_pages_path, output_path, position, verbose=False):
    original = Presentation(original_path)
    new_deck = Presentation(new_pages_path)

    if not (1 <= position <= len(original.slides) + 1):
        raise ValueError(
            f"--at {position} out of range (original has {len(original.slides)} slides; "
            f"valid positions: 1 to {len(original.slides) + 1})"
        )

    # Append first, then reorder
    for idx, new_slide in enumerate(new_deck.slides, 1):
        if verbose:
            print(f"  appending slide {idx}/{len(new_deck.slides)}", file=sys.stderr)
        copy_slide_into_deck(new_slide, original, verbose)

    if verbose:
        print(f"  moving {len(new_deck.slides)} appended slides to position {position}",
              file=sys.stderr)
    move_appended_to_position(original, len(new_deck.slides), position)

    original.save(output_path)
    return len(new_deck.slides)


def do_replace(original_path, new_pages_path, output_path, position, verbose=False):
    original = Presentation(original_path)
    new_deck = Presentation(new_pages_path)

    if not (1 <= position <= len(original.slides)):
        raise ValueError(
            f"--at {position} out of range (original has {len(original.slides)} slides)"
        )

    # Step 1: append new slides
    for idx, new_slide in enumerate(new_deck.slides, 1):
        if verbose:
            print(f"  appending slide {idx}/{len(new_deck.slides)}", file=sys.stderr)
        copy_slide_into_deck(new_slide, original, verbose)

    num_new = len(new_deck.slides)

    # Step 2: remove the new slides from end (we'll re-insert in step 4)
    slide_id_list = original.slides._sldIdLst
    all_ids = list(slide_id_list)
    new_ids = all_ids[len(all_ids) - num_new:]
    for sid in new_ids:
        slide_id_list.remove(sid)

    # Step 3: delete the target slide
    if verbose:
        print(f"  deleting slide {position}", file=sys.stderr)
    delete_slide_at(original, position)

    # Step 4: insert new slides at the target position
    insert_idx = position - 1
    for offset, sid in enumerate(new_ids):
        slide_id_list.insert(insert_idx + offset, sid)

    original.save(output_path)
    return num_new


# ============================================================
# Main
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Splice new pages from one PPTX into another.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__.split("Usage:", 1)[1] if "Usage:" in __doc__ else "",
    )
    parser.add_argument("--original", required=True,
                        help="Path to original .pptx (will not be modified)")
    parser.add_argument("--new-pages", required=True,
                        help="Path to .pptx containing the new pages to splice in")
    parser.add_argument("--mode", required=True, choices=["insert", "replace", "append"],
                        help="Splice mode")
    parser.add_argument("--at", type=int, default=None,
                        help="1-based position for insert/replace modes")
    parser.add_argument("--output", required=True,
                        help="Path for the output .pptx")
    parser.add_argument("--verbose", action="store_true",
                        help="Print progress details to stderr")

    args = parser.parse_args()

    original_path = Path(args.original)
    new_pages_path = Path(args.new_pages)
    output_path = Path(args.output)

    # Validate input files
    for label, p in [("--original", original_path), ("--new-pages", new_pages_path)]:
        if not p.exists():
            print(f"ERROR: {label} not found: {p}", file=sys.stderr)
            sys.exit(1)
        if p.suffix.lower() != ".pptx":
            print(f"ERROR: {label} is not a .pptx file: {p}", file=sys.stderr)
            sys.exit(1)

    # Validate --at for modes that need it
    if args.mode in ("insert", "replace") and args.at is None:
        print(f"ERROR: --at <position> is required for --mode {args.mode}",
              file=sys.stderr)
        sys.exit(1)

    # Ensure output directory exists
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Print summary
    if args.verbose:
        print(f"Original:  {original_path}", file=sys.stderr)
        print(f"New pages: {new_pages_path}", file=sys.stderr)
        print(f"Mode:      {args.mode}" + (f" at {args.at}" if args.at else ""),
              file=sys.stderr)
        print(f"Output:    {output_path}", file=sys.stderr)
        print(file=sys.stderr)

    # Dispatch
    try:
        if args.mode == "append":
            n = do_append(original_path, new_pages_path, output_path, args.verbose)
        elif args.mode == "insert":
            n = do_insert(original_path, new_pages_path, output_path, args.at,
                          args.verbose)
        elif args.mode == "replace":
            n = do_replace(original_path, new_pages_path, output_path, args.at,
                           args.verbose)
    except (ValueError, Exception) as e:
        print(f"ERROR: {e}", file=sys.stderr)
        sys.exit(1)

    # Final report
    final_prs = Presentation(output_path)
    print(f"OK. Spliced {n} new slide(s) into {output_path}")
    print(f"   Final slide count: {len(final_prs.slides)}")


if __name__ == "__main__":
    main()
