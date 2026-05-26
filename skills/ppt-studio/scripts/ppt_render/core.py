"""Core rendering primitives.

Coordinate system: 960 x 540 px canvas (16:9).
PPTX conversion: inches = px / 96.

An Element is a plain dict. layout() functions emit lists of these;
to_html() and to_pptx() consume them generically — no per-template code.
"""

from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Canvas constants ──
CANVAS_W = 960
CANVAS_H = 540
PX_PER_INCH = 96
EMU_PER_PX = 914400 // PX_PER_INCH  # 9525

# Hard upscale factor for screenshot crispness (HTML rendered at 2x)
PREVIEW_SCALE = 2


def px_to_emu(v):
    """Convert a pixel value to EMU for python-pptx."""
    return Emu(int(round(v * EMU_PER_PX)))


def hexclean(c):
    """Strip a leading # and uppercase a hex color."""
    if not c:
        return "000000"
    return c.lstrip("#").upper()


# ══════════════════════════════════════════════════════════
# Element factory helpers — layout() functions call these
# ══════════════════════════════════════════════════════════

def rect(x, y, w, h, fill, opacity=1.0, radius=0):
    return {"type": "rect", "x": x, "y": y, "w": w, "h": h,
            "fill": hexclean(fill), "opacity": opacity, "radius": radius}


def oval(x, y, w, h, fill, opacity=1.0):
    return {"type": "oval", "x": x, "y": y, "w": w, "h": h,
            "fill": hexclean(fill), "opacity": opacity}


def line(x, y, w, h, fill, opacity=1.0):
    """A line is just a thin rect (vertical: small w; horizontal: small h)."""
    return {"type": "rect", "x": x, "y": y, "w": w, "h": h,
            "fill": hexclean(fill), "opacity": opacity, "radius": 0}


def text(x, y, w, h, content, size=14, color="FFFFFF", font="Microsoft YaHei",
         bold=False, align="left", valign="top", wrap=True, line_height=1.4,
         opacity=1.0):
    return {"type": "text", "x": x, "y": y, "w": w, "h": h,
            "text": content if content is not None else "",
            "size": size, "color": hexclean(color), "font": font,
            "bold": bold, "align": align, "valign": valign,
            "wrap": wrap, "line_height": line_height, "opacity": opacity}


# ══════════════════════════════════════════════════════════
# HTML renderer
# ══════════════════════════════════════════════════════════

_ALIGN_CSS = {"left": "flex-start", "center": "center", "right": "flex-end"}
_VALIGN_CSS = {"top": "flex-start", "middle": "center", "bottom": "flex-end"}
_TEXTALIGN_CSS = {"left": "left", "center": "center", "right": "right"}


def _el_to_html(el):
    t = el["type"]
    if t in ("rect", "oval"):
        radius = "50%" if t == "oval" else f'{el.get("radius", 0)}px'
        op = el.get("opacity", 1.0)
        return (
            f'<div style="position:absolute;'
            f'left:{el["x"]}px;top:{el["y"]}px;'
            f'width:{el["w"]}px;height:{el["h"]}px;'
            f'background:#{el["fill"]};'
            f'opacity:{op};'
            f'border-radius:{radius};"></div>'
        )
    if t == "text":
        justify = _VALIGN_CSS.get(el.get("valign", "top"), "flex-start")
        align_items = _ALIGN_CSS.get(el.get("align", "left"), "flex-start")
        text_align = _TEXTALIGN_CSS.get(el.get("align", "left"), "left")
        weight = "bold" if el.get("bold") else "normal"
        wrap = "normal" if el.get("wrap", True) else "nowrap"
        op = el.get("opacity", 1.0)
        # \n in content becomes <br>
        safe = (el["text"]
                .replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))
        safe = safe.replace("\n", "<br>")
        return (
            f'<div style="position:absolute;'
            f'left:{el["x"]}px;top:{el["y"]}px;'
            f'width:{el["w"]}px;height:{el["h"]}px;'
            f'display:flex;flex-direction:column;'
            f'justify-content:{justify};align-items:{align_items};'
            f'overflow:hidden;opacity:{op};">'
            f'<span style="'
            f'font-family:\'{el["font"]}\',Arial,sans-serif;'
            f'font-size:{el["size"]}px;'
            f'font-weight:{weight};'
            f'color:#{el["color"]};'
            f'line-height:{el.get("line_height", 1.4)};'
            f'text-align:{text_align};'
            f'white-space:{wrap};'
            f'width:100%;">'
            f'{safe}</span></div>'
        )
    return ""


def slide_to_html(elements, label=None, scale=1):
    """Render one slide's elements into a 16:9 div."""
    inner = "".join(_el_to_html(e) for e in elements)
    w, h = CANVAS_W * scale, CANVAS_H * scale
    label_html = ""
    if label:
        label_html = (f'<div style="font-size:12px;color:#888;'
                      f'margin:0 0 6px 2px;font-family:Arial,sans-serif;">'
                      f'{label}</div>')
    transform = ""
    if scale != 1:
        transform = (f'transform:scale({scale});transform-origin:top left;')
    return (
        f'<div style="margin:0 auto 28px;width:{w}px;">'
        f'{label_html}'
        f'<div style="position:relative;width:{CANVAS_W}px;height:{CANVAS_H}px;'
        f'overflow:hidden;box-shadow:0 4px 24px rgba(0,0,0,0.18);'
        f'{transform}">'
        f'{inner}'
        f'</div></div>'
    )


def build_preview_html(slides_elements, labels=None):
    """slides_elements: list of (elements list). Returns full standalone HTML."""
    labels = labels or [None] * len(slides_elements)
    body = "".join(
        slide_to_html(els, label=labels[i])
        for i, els in enumerate(slides_elements)
    )
    return (
        '<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">'
        '<title>PPT Preview</title>'
        '<style>*{margin:0;padding:0;box-sizing:border-box;}'
        'body{background:#E8EAF0;padding:32px 16px;}</style>'
        '</head><body>'
        f'{body}'
        '</body></html>'
    )


def single_slide_html(elements):
    """A bare HTML page for one slide, exact canvas size — used for screenshots."""
    inner = "".join(_el_to_html(e) for e in elements)
    return (
        '<!DOCTYPE html><html><head><meta charset="UTF-8">'
        '<style>*{margin:0;padding:0;box-sizing:border-box;}'
        f'html,body{{width:{CANVAS_W}px;height:{CANVAS_H}px;overflow:hidden;}}'
        '</style></head><body>'
        f'<div style="position:relative;width:{CANVAS_W}px;height:{CANVAS_H}px;'
        f'overflow:hidden;">{inner}</div>'
        '</body></html>'
    )


# ══════════════════════════════════════════════════════════
# PPTX renderer
# ══════════════════════════════════════════════════════════

_PP_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
             "right": PP_ALIGN.RIGHT}
_MSO_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
               "bottom": MSO_ANCHOR.BOTTOM}


def _add_rect(slide, el, shape_type):
    shp = slide.shapes.add_shape(
        shape_type,
        px_to_emu(el["x"]), px_to_emu(el["y"]),
        px_to_emu(el["w"]), px_to_emu(el["h"]),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(el["fill"])
    shp.line.fill.background()
    shp.shadow.inherit = False
    _strip_shape_style(shp)
    op = el.get("opacity", 1.0)
    if op < 1.0:
        _set_shape_alpha(shp, op)
    return shp


def _strip_shape_style(shp):
    """Remove the <p:style> theme reference.

    python-pptx's add_shape() attaches a <p:style> with effectRef/fillRef
    pointing at theme styles. LibreOffice applies the effectRef (a drop
    shadow) even when an empty <a:effectLst/> is present. Removing the
    style element makes the explicit spPr the sole source of truth.
    """
    from pptx.oxml.ns import qn
    sp = shp._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)


def _set_corner_radius(shp, radius_px, shape_w_px):
    """Set rounded-rect corner radius. avLst val is a fraction of the
    shorter side, in 1/100000 units."""
    from pptx.oxml.ns import qn
    sp = shp._element
    geom = sp.find(qn("p:spPr") + "/" + qn("a:prstGeom"))
    if geom is None:
        return
    avlst = geom.find(qn("a:avLst"))
    if avlst is None:
        avlst = geom.makeelement(qn("a:avLst"), {})
        geom.append(avlst)
    for gd in avlst.findall(qn("a:gd")):
        avlst.remove(gd)
    # radius as fraction of the smaller dimension
    frac = min(0.5, radius_px / max(shape_w_px, 1))
    gd = avlst.makeelement(qn("a:gd"),
                           {"name": "adj", "fmla": f"val {int(frac * 100000)}"})
    avlst.append(gd)


def _set_shape_alpha(shp, opacity):
    """Inject an <a:alpha> element into the solidFill so the shape is translucent."""
    from pptx.oxml.ns import qn
    sp = shp.fill._xPr  # the <p:spPr> element
    srgb = sp.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
    # fallback search
    solid = sp.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # remove existing alpha
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
    srgb.append(alpha)


def _add_text(slide, el):
    from pptx.util import Pt
    box = slide.shapes.add_textbox(
        px_to_emu(el["x"]), px_to_emu(el["y"]),
        px_to_emu(el["w"]), px_to_emu(el["h"]),
    )
    tf = box.text_frame
    tf.word_wrap = el.get("wrap", True)
    tf.vertical_anchor = _MSO_ANCHOR.get(el.get("valign", "top"), MSO_ANCHOR.TOP)
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    opacity = el.get("opacity", 1.0)
    size = el["size"]
    lh = el.get("line_height", 1.4)
    # Point-based line spacing renders consistently across PowerPoint and
    # LibreOffice; the multiplier form is interpreted differently by each.
    spacing_pt = Pt(size * lh)
    lines = el["text"].split("\n")
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = _PP_ALIGN.get(el.get("align", "left"), PP_ALIGN.LEFT)
        try:
            para.line_spacing = spacing_pt
        except Exception:
            pass
        run = para.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.name = el["font"]
        run.font.bold = el.get("bold", False)
        run.font.color.rgb = RGBColor.from_string(el["color"])
        _set_ea_font(run, el["font"])
        if opacity < 1.0:
            _set_run_alpha(run, opacity)
    return box


def _set_run_alpha(run, opacity):
    """Apply alpha to a run's color (for faint decorative text)."""
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    fill = rPr.find(qn("a:solidFill"))
    if fill is None:
        return
    srgb = fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for a in srgb.findall(qn("a:alpha")):
        srgb.remove(a)
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(opacity * 100000))})
    srgb.append(alpha)


def _set_ea_font(run, font_name):
    """Ensure CJK text uses the intended font (python-pptx only sets latin)."""
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def render_slide_to_pptx(slide, elements):
    """Draw all elements onto a python-pptx slide."""
    for el in elements:
        t = el["type"]
        if t == "rect":
            radius = el.get("radius", 0)
            if radius > 0:
                shp = _add_rect(slide, el, MSO_SHAPE.ROUNDED_RECTANGLE)
                _set_corner_radius(shp, radius, el["w"])
            else:
                _add_rect(slide, el, MSO_SHAPE.RECTANGLE)
        elif t == "oval":
            _add_rect(slide, el, MSO_SHAPE.OVAL)
        elif t == "text":
            _add_text(slide, el)
