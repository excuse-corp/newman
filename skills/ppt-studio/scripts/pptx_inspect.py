#!/usr/bin/env python3
"""
pptx_inspect.py — list all slides and shapes in a PPTX for E-path editing.

Usage:
    python pptx_inspect.py <file.pptx>              # human-readable text
    python pptx_inspect.py <file.pptx> --json       # machine-readable JSON
    python pptx_inspect.py <file.pptx> --slide 3    # only slide 3 (1-based)

Output format (text mode):
    === Slide 1 ===
      [s1_0]  TEXT_BOX     "title text..."           x=0.42 y=2.10 w=5.80 h=0.80  runs=1
      [s1_3]  PICTURE      <image, 320x240>           x=6.50 y=1.20 w=2.80 h=2.10
      [s1_4]  GROUP        (3 shapes inside)          x=0.60 y=1.20 w=8.80 h=3.20
        [s1_4.0]  AUTO_SHAPE  "" (fill=#3D2A1E)       x=0.60 y=1.20 w=2.80 h=3.20

Shape ID convention:
    s{slide}_{shape}            top-level shape (1-based slide, 0-based shape)
    s{slide}_{group}.{child}    shape inside a group
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu


# ---------- helpers ----------

def emu_to_inches(emu):
    """Convert EMU to inches. Returns None if input is None."""
    if emu is None:
        return None
    return round(Emu(emu).inches, 2)


def display_width(s):
    """Approximate display width: CJK and full-width chars count as 2, others 1."""
    import unicodedata
    width = 0
    for ch in s:
        # 'W' = wide, 'F' = full-width, 'A' = ambiguous (treat as 1 here)
        if unicodedata.east_asian_width(ch) in ("W", "F"):
            width += 2
        else:
            width += 1
    return width


def lpad(s, target_width):
    """Pad string with spaces on the right to reach target display width."""
    current = display_width(s)
    if current >= target_width:
        return s
    return s + " " * (target_width - current)


def shape_type_name(shape):
    """Return a short, readable shape type name."""
    st = shape.shape_type
    if st is None:
        return "UNKNOWN"
    
    type_map = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AUTO_SHAPE",
        MSO_SHAPE_TYPE.TEXT_BOX: "TEXT_BOX",
        MSO_SHAPE_TYPE.PICTURE: "PICTURE",
        MSO_SHAPE_TYPE.PLACEHOLDER: "PLACEHOLDER",
        MSO_SHAPE_TYPE.GROUP: "GROUP",
        MSO_SHAPE_TYPE.TABLE: "TABLE",
        MSO_SHAPE_TYPE.CHART: "CHART",
        MSO_SHAPE_TYPE.LINE: "LINE",
        MSO_SHAPE_TYPE.FREEFORM: "FREEFORM",
        MSO_SHAPE_TYPE.MEDIA: "MEDIA",
        MSO_SHAPE_TYPE.DIAGRAM: "SMART_ART",
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "OLE_OBJECT",
        MSO_SHAPE_TYPE.LINKED_PICTURE: "LINKED_PICTURE",
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "LINKED_OLE",
    }
    return type_map.get(st, f"TYPE_{int(st)}")


def text_preview(shape, max_chars=40):
    """Extract a short text preview from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    text = shape.text_frame.text.strip().replace("\n", " ").replace("\r", " ")
    if not text:
        return ""
    if len(text) > max_chars:
        return text[:max_chars] + "…"
    return text


def count_runs(shape):
    """Count total runs across all paragraphs in a text frame."""
    if not shape.has_text_frame:
        return 0
    return sum(len(p.runs) for p in shape.text_frame.paragraphs)


def get_fill_hex(shape):
    """Try to extract solid fill color as #RRGGBB hex. Returns None if not applicable."""
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        # type 1 = MSO_FILL_TYPE.SOLID
        if fill.type != 1:
            return None
        rgb = fill.fore_color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb).upper()}"
    except (AttributeError, ValueError, Exception):
        return None


def get_picture_dimensions(shape):
    """For PICTURE shapes, get the embedded image's pixel dimensions."""
    try:
        img = shape.image
        return f"{img.size[0]}x{img.size[1]}"
    except (AttributeError, ValueError):
        return "unknown"


def get_table_dimensions(shape):
    """For TABLE shapes, return rows x cols."""
    try:
        tbl = shape.table
        rows = len(tbl.rows)
        cols = len(tbl.columns)
        return f"{rows}x{cols}"
    except (AttributeError, ValueError):
        return "unknown"


def is_placeholder(shape):
    """Check whether shape is a placeholder (with safe attribute access)."""
    try:
        return shape.is_placeholder
    except (AttributeError, ValueError):
        return False


def placeholder_info(shape):
    """For placeholders, return type and idx."""
    try:
        ph = shape.placeholder_format
        return f"ph_idx={ph.idx}, ph_type={ph.type}"
    except (AttributeError, ValueError):
        return ""


# ---------- shape inspection ----------

def inspect_shape(shape, shape_id):
    """Return a dict describing a single shape (recursive for groups)."""
    info = {
        "id": shape_id,
        "type": shape_type_name(shape),
        "name": getattr(shape, "name", ""),
        "x": emu_to_inches(shape.left),
        "y": emu_to_inches(shape.top),
        "w": emu_to_inches(shape.width),
        "h": emu_to_inches(shape.height),
        "text": text_preview(shape),
        "runs": count_runs(shape),
        "is_placeholder": is_placeholder(shape),
    }
    
    # Placeholder details
    if info["is_placeholder"]:
        info["placeholder"] = placeholder_info(shape)
    
    # Fill color (for AUTO_SHAPE / TEXT_BOX with explicit fill)
    fill = get_fill_hex(shape)
    if fill:
        info["fill"] = fill
    
    # Type-specific extras
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["image_size"] = get_picture_dimensions(shape)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        info["table_size"] = get_table_dimensions(shape)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = []
        for child_idx, child in enumerate(shape.shapes):
            child_id = f"{shape_id}.{child_idx}"
            info["children"].append(inspect_shape(child, child_id))
    
    return info


def inspect_slide(slide, slide_num):
    """Return a dict describing one slide and all its shapes."""
    shapes_info = []
    for shape_idx, shape in enumerate(slide.shapes):
        shape_id = f"s{slide_num}_{shape_idx}"
        shapes_info.append(inspect_shape(shape, shape_id))
    
    return {
        "slide_num": slide_num,
        "layout": slide.slide_layout.name if slide.slide_layout else "",
        "shapes": shapes_info,
    }


def inspect_presentation(pptx_path, slide_filter=None):
    """Inspect all slides (or just one if slide_filter given, 1-based)."""
    prs = Presentation(pptx_path)
    
    deck_info = {
        "file": str(pptx_path),
        "slide_count": len(prs.slides),
        "slide_width_in": emu_to_inches(prs.slide_width),
        "slide_height_in": emu_to_inches(prs.slide_height),
        "slides": [],
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1  # 1-based
        if slide_filter is not None and slide_num != slide_filter:
            continue
        deck_info["slides"].append(inspect_slide(slide, slide_num))
    
    return deck_info


# ---------- text formatting ----------

def format_shape_line(info, indent=2):
    """Format a single shape's info as a one-line text summary."""
    pad = " " * indent
    
    # Shape ID and type
    id_part = f"[{info['id']}]"
    type_part = info["type"]
    
    # Content preview (text / image-size / table-size / group-children-count)
    if info["type"] == "PICTURE":
        content = f"<image, {info.get('image_size', '?')}>"
    elif info["type"] == "TABLE":
        content = f"<table, {info.get('table_size', '?')}>"
    elif info["type"] == "GROUP":
        n = len(info.get("children", []))
        content = f"(GROUP, {n} shapes inside)"
    elif info["type"] == "CHART":
        content = "<chart, opaque to E-path>"
    elif info["type"] == "SMART_ART":
        content = "<smart_art, opaque to E-path>"
    elif info["text"]:
        content = f'"{info["text"]}"'
    else:
        content = '""'
    
    # Fill annotation
    if info.get("fill"):
        content += f" (fill={info['fill']})"
    
    # Placeholder annotation
    if info.get("is_placeholder") and info.get("placeholder"):
        content += f" [{info['placeholder']}]"
    
    # Position
    x, y, w, h = info["x"], info["y"], info["w"], info["h"]
    pos = f"x={x} y={y} w={w} h={h}" if all(v is not None for v in [x, y, w, h]) else ""
    
    # Runs (only for text-bearing shapes)
    runs_part = f"runs={info['runs']}" if info["runs"] > 0 else ""
    
    line = f"{pad}{lpad(id_part, 12)} {lpad(type_part, 12)} {content}"
    # Pad to a target column for alignment, then position info
    if pos:
        line = lpad(line, 70) + "  " + pos
    if runs_part:
        line = f"{line}  {runs_part}"
    
    return line


def print_shape_tree(info, indent=2):
    """Print a shape and recursively its children (for groups)."""
    print(format_shape_line(info, indent))
    if info["type"] == "GROUP":
        for child in info.get("children", []):
            print_shape_tree(child, indent + 2)


def print_text_report(deck_info):
    """Print human-readable inspection report."""
    print(f"File: {deck_info['file']}")
    print(f"Slides: {deck_info['slide_count']}  "
          f"Page size: {deck_info['slide_width_in']}\" x {deck_info['slide_height_in']}\"")
    print()
    
    for slide in deck_info["slides"]:
        print(f"=== Slide {slide['slide_num']} ===  layout: {slide['layout']}")
        for shape_info in slide["shapes"]:
            print_shape_tree(shape_info)
        print()


def print_json_report(deck_info):
    """Print machine-readable JSON."""
    print(json.dumps(deck_info, ensure_ascii=False, indent=2))


# ---------- main ----------

def main():
    parser = argparse.ArgumentParser(
        description="Inspect a PPTX file: list all slides and shapes with IDs for E-path editing.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Shape IDs use format s{slide_1based}_{shape_0based}, e.g. s3_2 = slide 3, shape index 2.",
    )
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("--json", action="store_true", help="Output JSON instead of text")
    parser.add_argument("--slide", type=int, default=None,
                        help="Only inspect this slide (1-based)")
    
    args = parser.parse_args()
    
    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"ERROR: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    if pptx_path.suffix.lower() != ".pptx":
        print(f"ERROR: not a .pptx file: {pptx_path}", file=sys.stderr)
        print("If this is a .ppt (legacy binary), open it in PowerPoint and Save As .pptx first.",
              file=sys.stderr)
        sys.exit(1)
    
    try:
        deck_info = inspect_presentation(pptx_path, slide_filter=args.slide)
    except Exception as e:
        print(f"ERROR: failed to open {pptx_path}: {e}", file=sys.stderr)
        sys.exit(1)
    
    # Check if --slide was given but produced no results
    if args.slide is not None and not deck_info["slides"]:
        print(f"ERROR: slide {args.slide} is out of range. "
              f"Deck has {deck_info['slide_count']} slides (1-{deck_info['slide_count']}).",
              file=sys.stderr)
        sys.exit(1)
    
    if args.json:
        print_json_report(deck_info)
    else:
        print_text_report(deck_info)


if __name__ == "__main__":
    main()
