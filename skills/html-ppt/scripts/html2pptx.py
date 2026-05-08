#!/usr/bin/env python3
"""
html2pptx.py — Hybrid PPT generator (Python rewrite of html2pptx.js)

Renders slide backgrounds as HTML → PNG screenshots (CSS gradients, shadows,
rounded corners), then overlays editable text via python-pptx.

Usage: python3 html2pptx.py <input.json> [output.pptx]
"""

import sys
import json
import base64
import tempfile
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from playwright.sync_api import sync_playwright

# ── Slide canvas constants ──
W, H = 1920, 1080          # pixel dimensions for HTML rendering
IW, IH = 10.0, 5.625      # PPTX slide size in inches

def px(v, axis):
    """Convert pixel value to inches for PPTX."""
    return (v / W) * IW if axis == "x" else (v / H) * IH

def thC(th, key):
    """Return '#RRGGBB' from theme dict."""
    val = th.get(key, "FFFFFF")
    return "#" + val if not val.startswith("#") else val

def rgb(hex_str):
    """Parse hex color string (with or without #) to RGBColor."""
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    h = h[:6].ljust(6, "0")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ══════════════════════════════════════════════════════════
# CSS RESET — shared by all slide HTML templates
# ══════════════════════════════════════════════════════════
RESET = f"* {{ margin:0; padding:0; box-sizing:border-box; }}\nbody {{ width:{W}px; height:{H}px; overflow:hidden; }}"

def vis(show):
    return "visible" if show else "hidden"

# ── Shared HTML helpers ──
def title_bar_html(th, title, show):
    s = vis(show)
    return f"""<div style="position:absolute;left:0;top:0;right:0;height:130px;
    background:linear-gradient(135deg,{thC(th,'primary')},{thC(th,'primary')}dd);
    display:flex;align-items:center;padding:0 80px;">
    <div style="width:5px;height:40px;background:{thC(th,'secondary')};margin-right:20px;border-radius:2px;"></div>
    <span style="font-size:36px;font-weight:bold;color:{thC(th,'accent')};font-family:Georgia,serif;visibility:{s};">{title or ''}</span>
  </div>"""

def page_badge_html(th, num, show):
    s = vis(show)
    return f"""<div style="position:absolute;right:40px;bottom:28px;width:38px;height:38px;border-radius:50%;
    background:{thC(th,'secondary')};display:flex;align-items:center;justify-content:center;
    font-size:14px;color:white;font-weight:bold;visibility:{s};">{num}</div>"""

def card_css(th):
    return "background:white;border-radius:12px;padding:32px;position:relative;overflow:hidden;box-shadow:0 2px 16px rgba(0,0,0,0.07);border:1px solid #e8e8e8;"

def card_top_accent(th):
    return f"content:'';position:absolute;top:0;left:0;right:0;height:4px;background:{thC(th,'secondary')};border-radius:12px 12px 0 0;"

# ══════════════════════════════════════════════════════════
# TEXT OVERLAY HELPER
# ══════════════════════════════════════════════════════════
def T(x, y, w, h, text, fs=14, c="FFFFFF", ff="Calibri", b=False, a="left", va="middle", i=False):
    """Build a text overlay descriptor dict."""
    return dict(x=x, y=y, w=w, h=h, text=text, fs=fs, c=c, ff=ff, b=b, a=a, va=va, italic=i)

# ══════════════════════════════════════════════════════════
# TEMPLATE REGISTRY
# Each template has:
#   html(th, d, show) → HTML string for background rendering
#   texts(th, d)      → list of T() dicts for text overlay
# ══════════════════════════════════════════════════════════
TEMPLATES = {}

# ── 1. cover_pro ──
def cover_pro_html(th, d, show):
    s = vis(show)
    stats = "".join(
        f"""<div style="display:flex;align-items:center;gap:12px;">
        <div style="width:10px;height:10px;border-radius:50%;background:{thC(th,'secondary')};"></div>
        <span style="font-size:17px;color:{thC(th,'secondary')};letter-spacing:1px;visibility:{s};">{st.get('text','')}</span>
      </div>"""
        for st in (d.get("bottom_stats") or [])
    )
    title_html = (d.get("title") or "").replace("\n", "<br>")
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Calibri,Arial,sans-serif; }}
    </style></head><body>
      <div style="position:absolute;left:0;top:0;width:42%;height:100%;background:linear-gradient(135deg,{thC(th,'secondary')}15,{thC(th,'secondary')}30);"></div>
      <div style="position:absolute;right:-60px;top:-60px;width:360px;height:360px;border-radius:50%;border:2px solid {thC(th,'secondary')}33;"></div>
      <div style="position:absolute;right:-20px;top:-20px;width:280px;height:280px;border-radius:50%;border:1px solid {thC(th,'secondary')}1a;"></div>
      <div style="position:absolute;left:80px;top:38px;font-size:15px;color:{thC(th,'secondary')};letter-spacing:3px;text-transform:uppercase;border-left:3px solid {thC(th,'secondary')};padding-left:12px;visibility:{s};">{d.get('department','')}</div>
      <div style="position:absolute;left:80px;top:180px;width:85%;font-size:72px;font-weight:bold;color:{thC(th,'accent')};font-family:Georgia,serif;line-height:1.15;visibility:{s};">{title_html}</div>
      <div style="position:absolute;left:80px;top:510px;display:flex;align-items:center;gap:20px;visibility:{s};">
        <div style="width:60px;height:2px;background:{thC(th,'secondary')};"></div>
        <span style="font-size:26px;color:{thC(th,'secondary')};">{d.get('subtitle','')}</span>
      </div>
      <div style="position:absolute;bottom:0;left:0;right:0;height:90px;border-top:1px solid {thC(th,'secondary')}33;display:flex;align-items:center;padding:0 80px;gap:50px;">
        {stats}
        <span style="margin-left:auto;font-size:17px;color:{thC(th,'secondary')};letter-spacing:2px;visibility:{s};">{d.get('date','')}</span>
      </div>
    </body></html>"""

def cover_pro_texts(th, d):
    t = [
        T(80, 38, 600, 28, d.get("department",""), fs=11, c=th.get("secondary","FFFFFF")),
        T(80, 180, 1600, 300, (d.get("title","") or "").replace("\n","\n"), fs=54, c=th.get("accent","FFFFFF"), ff="Georgia", b=True),
        T(160, 510, 1400, 45, d.get("subtitle",""), fs=20, c=th.get("secondary","FFFFFF")),
        T(1600, 1000, 250, 30, d.get("date",""), fs=13, c=th.get("secondary","FFFFFF"), a="right"),
    ]
    for i, st in enumerate(d.get("bottom_stats") or []):
        t.append(T(100+i*280, 1000, 250, 30, st.get("text",""), fs=12, c=th.get("secondary","FFFFFF")))
    return t

TEMPLATES["cover_pro"] = {"html": cover_pro_html, "texts": cover_pro_texts}
TEMPLATES["cover"] = TEMPLATES["cover_pro"]

# ── 2. section_break ──
def section_break_html(th, d, show):
    s = vis(show)
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Georgia,serif; }}
    </style></head><body>
      <div style="position:absolute;left:-40px;top:80px;font-size:280px;font-weight:bold;color:{thC(th,'secondary')}18;line-height:1;visibility:{s};">{d.get('section_number','')}</div>
      <div style="position:absolute;right:50px;top:12%;width:4px;height:70%;background:{thC(th,'secondary')}44;border-radius:2px;"></div>
      <div style="position:absolute;left:80px;top:420px;font-size:44px;font-weight:bold;color:{thC(th,'accent')};visibility:{s};">{d.get('title','')}</div>
      <div style="position:absolute;left:80px;top:500px;font-size:20px;color:{thC(th,'secondary')};visibility:{s};">{d.get('subtitle','')}</div>
    </body></html>"""

def section_break_texts(th, d):
    return [
        T(0, 80, 700, 300, d.get("section_number",""), fs=180, c=th.get("secondary","FFFFFF")+"30", ff="Georgia", b=True, a="center"),
        T(80, 420, 1400, 60, d.get("title",""), fs=36, c=th.get("accent","FFFFFF"), ff="Georgia", b=True),
        T(80, 500, 1400, 40, d.get("subtitle",""), fs=16, c=th.get("secondary","FFFFFF")),
    ]

TEMPLATES["section_break"] = {"html": section_break_html, "texts": section_break_texts}

# ── 3. card_grid ──
def card_grid_html(th, d, show):
    s = vis(show)
    cards = d.get("cards") or []
    cols = d.get("columns") or 3
    cards_html = "".join(
        f"""<div class="card"><div class="icon"><div class="dot"></div></div>
          <div style="font-size:20px;font-weight:bold;color:#1a1a2e;margin-bottom:8px;visibility:{s};">{c.get('heading','')}</div>
          <div style="font-size:14px;color:#666;line-height:1.65;visibility:{s};">{c.get('text','')}</div>
        </div>"""
        for c in cards
    )
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:#F5F5F5; font-family:Calibri,Arial,sans-serif; }}
      .card {{ {card_css(th)} }}
      .card::before {{ {card_top_accent(th)} }}
      .icon {{ width:44px;height:44px;border-radius:10px;background:{thC(th,'secondary')}18;display:flex;align-items:center;justify-content:center;margin-bottom:14px; }}
      .dot {{ width:18px;height:18px;border-radius:50%;background:{thC(th,'secondary')}; }}
    </style></head><body>
      {title_bar_html(th, d.get('title'), show)}
      <div style="position:absolute;left:60px;top:160px;right:60px;bottom:36px;display:grid;grid-template-columns:repeat({cols},1fr);gap:22px;">
        {cards_html}
      </div>
      {page_badge_html(th, d.get('page'), show)}
    </body></html>"""

def card_grid_texts(th, d):
    cards = d.get("cards") or []
    cols = d.get("columns") or 3
    gL, gT, gR, gB, gap = 60, 160, 60, 36, 22
    gW_px = W - gL - gR
    gH_px = H - gT - gB
    rows = -(-len(cards) // cols)  # ceil division
    cW = (gW_px - (cols-1)*gap) / cols
    cH = (gH_px - (rows-1)*gap) / rows if rows > 0 else gH_px
    t = [T(110, 36, 1500, 58, d.get("title",""), fs=28, c=th.get("accent","FFFFFF"), ff="Georgia", b=True)]
    for i, c in enumerate(cards):
        col = i % cols
        row = i // cols
        cx = gL + col*(cW+gap)
        cy = gT + row*(cH+gap)
        t.append(T(cx+32, cy+70, cW-64, 30, c.get("heading",""), fs=16, c="1A1A2E", b=True))
        t.append(T(cx+32, cy+108, cW-64, cH-140, c.get("text",""), fs=11, c="666666", va="top"))
    return t

TEMPLATES["card_grid"] = {"html": card_grid_html, "texts": card_grid_texts}

# ── 4. kpi_dashboard ──
def kpi_dashboard_html(th, d, show):
    s = vis(show)
    kpis = d.get("kpis") or []
    details = d.get("detail_cards") or []
    kpis_html = "".join(
        f"""<div style="flex:1;background:{thC(th,'primary')};border-radius:10px;display:flex;flex-direction:column;align-items:center;justify-content:center;">
          <span style="font-size:42px;font-weight:bold;color:{thC(th,'accent')};font-family:Georgia,serif;visibility:{s};">{k.get('value','')}</span>
          <span style="font-size:13px;color:{thC(th,'secondary')};margin-top:4px;visibility:{s};">{k.get('label','')}</span>
        </div>"""
        for k in kpis
    )
    details_html = "".join(
        f"""<div class="card" style="flex:1;">
          <div style="font-size:18px;font-weight:bold;color:#1a1a2e;margin-bottom:10px;visibility:{s};">{det.get('heading','')}</div>
          <div style="font-size:14px;color:#666;line-height:1.65;visibility:{s};">{det.get('text','')}</div>
        </div>"""
        for det in details
    )
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:#F5F5F5; font-family:Calibri,Arial,sans-serif; }}
      .card {{ {card_css(th)} }}
      .card::before {{ {card_top_accent(th)} }}
    </style></head><body>
      {title_bar_html(th, d.get('title'), show)}
      <div style="position:absolute;left:50px;top:160px;right:50px;height:120px;display:flex;gap:16px;">
        {kpis_html}
      </div>
      <div style="position:absolute;left:50px;top:300px;right:50px;bottom:36px;display:flex;gap:20px;">
        {details_html}
      </div>
      {page_badge_html(th, d.get('page'), show)}
    </body></html>"""

def kpi_dashboard_texts(th, d):
    kpis = d.get("kpis") or []
    dets = d.get("detail_cards") or []
    kC, dC = len(kpis), len(dets)
    kGap, kL, kR = 16, 50, 50
    kW = (W - kL - kR - (kC-1)*kGap) / kC if kC else W
    t = [T(110, 36, 1500, 58, d.get("title",""), fs=28, c=th.get("accent","FFFFFF"), ff="Georgia", b=True)]
    for i, k in enumerate(kpis):
        x = kL + i*(kW+kGap)
        t.append(T(x, 165, kW, 65, k.get("value",""), fs=34, c=th.get("accent","FFFFFF"), ff="Georgia", b=True, a="center"))
        t.append(T(x, 238, kW, 25, k.get("label",""), fs=11, c=th.get("secondary","FFFFFF"), a="center"))
    dGap, dL, dR = 20, 50, 50
    dW = (W - dL - dR - (dC-1)*dGap) / dC if dC else W
    for i, det in enumerate(dets):
        x = dL + i*(dW+dGap)
        t.append(T(x+32, 320, dW-64, 30, det.get("heading",""), fs=15, c="1A1A2E", b=True))
        t.append(T(x+32, 358, dW-64, 600, det.get("text",""), fs=11, c="666666", va="top"))
    return t

TEMPLATES["kpi_dashboard"] = {"html": kpi_dashboard_html, "texts": kpi_dashboard_texts}

# ── 5. grid_content ──
def grid_content_html(th, d, show):
    s = vis(show)
    lc = d.get("left_card") or {}
    rcs = d.get("right_cards") or []
    bb = d.get("bottom_bar")
    bb_bottom = "110" if bb else "36"
    right_cards_html = "".join(
        f"""<div class="card" style="flex:1;padding:20px;">
            <div style="font-size:16px;font-weight:bold;color:{thC(th,'text_dark')};margin-bottom:12px;visibility:{s};">{rc.get('heading','')}</div>
            <div style="display:flex;gap:16px;">
              {"".join(f'''<div style="flex:1;">
                <div style="font-size:13px;font-weight:bold;color:{thC(th,'secondary')};margin-bottom:4px;visibility:{s};">{it.get('label','')}</div>
                <div style="font-size:12px;color:{thC(th,'text_body')};line-height:1.6;visibility:{s};">{it.get('text','')}</div>
              </div>''' for it in (rc.get('items') or []))}
            </div>
          </div>"""
        for rc in rcs
    )
    bb_html = f"""<div style="position:absolute;left:40px;bottom:20px;right:40px;height:76px;background:{thC(th,'secondary')};border-radius:6px;padding:14px 24px;">
        <div style="font-size:15px;font-weight:bold;color:{thC(th,'primary')};visibility:{s};">{bb.get('heading','')}</div>
        <div style="font-size:12px;color:{thC(th,'primary')}cc;margin-top:4px;visibility:{s};">{bb.get('text','')}</div>
      </div>""" if bb else ""
    subtitle_html = f"""<div style="position:absolute;left:65px;top:80px;font-size:14px;color:{thC(th,'text_body')};visibility:{s};">{d.get('subtitle','')}</div>""" if d.get("subtitle") else ""
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Calibri,Arial,sans-serif; }}
      .card {{ background:{thC(th,'bg_light')};border-radius:8px;border:1px solid {thC(th,'secondary')}33;overflow:hidden; }}
    </style></head><body>
      <div style="position:absolute;left:40px;top:28px;display:flex;align-items:center;gap:16px;">
        <div style="width:5px;height:44px;background:{thC(th,'secondary')};border-radius:2px;"></div>
        <span style="font-size:30px;font-weight:bold;color:{thC(th,'accent')};font-family:Georgia,serif;visibility:{s};">{d.get('title','')}</span>
      </div>
      {subtitle_html}
      <div class="card" style="position:absolute;left:40px;top:115px;width:700px;bottom:{bb_bottom}px;padding:24px;">
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:16px;">
          <div style="background:{thC(th,'secondary')};color:{thC(th,'primary')};font-size:18px;font-weight:bold;padding:6px 14px;border-radius:4px;visibility:{s};">{lc.get('number','01')}</div>
          <span style="font-size:18px;font-weight:bold;color:{thC(th,'text_dark')};visibility:{s};">{lc.get('heading','')}</span>
        </div>
        <div style="font-size:14px;color:{thC(th,'text_body')};line-height:1.7;white-space:pre-wrap;visibility:{s};">{lc.get('content','')}</div>
      </div>
      <div style="position:absolute;left:760px;top:115px;right:40px;bottom:{bb_bottom}px;display:flex;flex-direction:column;gap:14px;">
        {right_cards_html}
      </div>
      {bb_html}
      {page_badge_html(th, d.get('page'), show)}
    </body></html>"""

def grid_content_texts(th, d):
    lc = d.get("left_card") or {}
    rcs = d.get("right_cards") or []
    bb = d.get("bottom_bar")
    t = [T(65, 28, 1500, 50, d.get("title",""), fs=24, c=th.get("accent","FFFFFF"), ff="Georgia", b=True)]
    if d.get("subtitle"):
        t.append(T(65, 80, 1500, 24, d["subtitle"], fs=11, c=th.get("text_body","FFFFFF")))
    t.append(T(80, 120, 80, 35, lc.get("number","01"), fs=15, c=th.get("primary","000000"), b=True, a="center"))
    t.append(T(170, 120, 500, 35, lc.get("heading",""), fs=15, c=th.get("text_dark","000000"), b=True))
    t.append(T(65, 175, 650, 600, (lc.get("content","") or "").replace("\\n","\n"), fs=11, c=th.get("text_body","FFFFFF"), va="top"))
    bbH = 110 if bb else 36
    rcTotalH = H - 115 - bbH
    rcH = (rcTotalH - (len(rcs)-1)*14) / len(rcs) if rcs else rcTotalH
    for i, rc in enumerate(rcs):
        ry = 115 + i*(rcH+14)
        t.append(T(785, ry+20, 800, 24, rc.get("heading",""), fs=13, c=th.get("text_dark","000000"), b=True))
        for j, it in enumerate(rc.get("items") or []):
            ix = 785 + j*520
            t.append(T(ix, ry+50, 490, 20, it.get("label",""), fs=10, c=th.get("secondary","FFFFFF"), b=True))
            t.append(T(ix, ry+74, 490, rcH-100, it.get("text",""), fs=10, c=th.get("text_body","FFFFFF"), va="top"))
    if bb:
        t.append(T(65, H-92, 800, 24, bb.get("heading",""), fs=12, c=th.get("primary","000000"), b=True))
        t.append(T(65, H-68, 1700, 22, bb.get("text",""), fs=10, c=th.get("primary","000000")))
    return t

TEMPLATES["grid_content"] = {"html": grid_content_html, "texts": grid_content_texts}

# ── 6. structured_content ──
def structured_content_html(th, d, show):
    s = vis(show)
    secs = d.get("sections") or []
    bb = d.get("bottom_bar")
    bb_bottom = "105" if bb else "30"
    subtitle_html = f"""<div style="position:absolute;left:65px;top:80px;font-size:14px;color:{thC(th,'text_body')};visibility:{s};">{d.get('subtitle','')}</div>""" if d.get("subtitle") else ""
    secs_html = "".join(
        f"""{'<div style="margin:0 24px;border-top:1px dashed ' + thC(th,'secondary') + '33;"></div>' if i > 0 else ''}
          <div style="padding:14px 24px;">
            <div style="font-size:15px;font-weight:bold;color:{thC(th,'accent')};margin-bottom:4px;visibility:{s};">{sec.get('heading','')}</div>
            <div style="font-size:13px;color:{thC(th,'text_body')};line-height:1.6;visibility:{s};">{sec.get('text','')}</div>
          </div>"""
        for i, sec in enumerate(secs)
    )
    callout_html = f"""<div style="margin:8px 24px;padding:10px 20px;background:{thC(th,'bg_light')};border-radius:6px;border:1px solid {thC(th,'secondary')}22;text-align:center;font-size:13px;color:{thC(th,'text_dark')};visibility:{s};">{d['callout'].get('text','')}</div>""" if d.get("callout") else ""
    bb_html = f"""<div style="position:absolute;left:40px;bottom:18px;right:40px;height:70px;background:{thC(th,'secondary')};border-radius:6px;padding:12px 24px;">
        <div style="font-size:14px;font-weight:bold;color:{thC(th,'primary')};visibility:{s};">{(bb.get('heading','')+'：' if bb.get('heading') else '')+bb.get('text','')}</div>
      </div>""" if bb else ""
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Calibri,Arial,sans-serif; }}
    </style></head><body>
      <div style="position:absolute;left:40px;top:28px;display:flex;align-items:center;gap:16px;">
        <div style="width:5px;height:44px;background:{thC(th,'secondary')};border-radius:2px;"></div>
        <span style="font-size:30px;font-weight:bold;color:{thC(th,'accent')};font-family:Georgia,serif;visibility:{s};">{d.get('title','')}</span>
      </div>
      {subtitle_html}
      <div style="position:absolute;left:40px;top:110px;right:40px;bottom:{bb_bottom}px;border:1px solid {thC(th,'secondary')}44;border-radius:8px;overflow:hidden;">
        <div style="display:flex;align-items:center;gap:12px;padding:18px 24px;border-bottom:1px solid {thC(th,'secondary')}22;">
          <div style="background:{thC(th,'secondary')};color:{thC(th,'primary')};font-size:16px;font-weight:bold;padding:5px 13px;border-radius:4px;visibility:{s};">{d.get('number_badge','01')}</div>
          <span style="font-size:18px;font-weight:bold;color:{thC(th,'accent')};visibility:{s};">{d.get('heading','')}</span>
        </div>
        {secs_html}
        {callout_html}
      </div>
      {bb_html}
      {page_badge_html(th, d.get('page'), show)}
    </body></html>"""

def structured_content_texts(th, d):
    secs = d.get("sections") or []
    bb = d.get("bottom_bar")
    t = [T(65, 28, 1500, 50, d.get("title",""), fs=24, c=th.get("accent","FFFFFF"), ff="Georgia", b=True)]
    if d.get("subtitle"):
        t.append(T(65, 80, 1500, 24, d["subtitle"], fs=11, c=th.get("text_body","FFFFFF")))
    t.append(T(85, 118, 80, 30, d.get("number_badge","01"), fs=13, c=th.get("primary","000000"), b=True, a="center"))
    t.append(T(175, 118, 1500, 30, d.get("heading",""), fs=15, c=th.get("accent","FFFFFF"), b=True))
    cardTop = 165
    cardBot = 975 if bb else H-30
    secH = (cardBot-cardTop-60) / max(len(secs), 1)
    for i, sec in enumerate(secs):
        sy = cardTop + i*secH
        t.append(T(65, sy+5, 1700, 24, sec.get("heading",""), fs=12, c=th.get("accent","FFFFFF"), b=True))
        t.append(T(65, sy+30, 1700, secH-40, sec.get("text",""), fs=10, c=th.get("text_body","FFFFFF"), va="top"))
    if d.get("callout"):
        t.append(T(100, cardBot-45, 1700, 30, d["callout"].get("text",""), fs=10, c=th.get("text_dark","000000"), a="center"))
    if bb:
        t.append(T(65, H-82, 1700, 30, (bb.get("heading","")+"：" if bb.get("heading") else "")+bb.get("text",""), fs=11, c=th.get("primary","000000"), b=True))
    return t

TEMPLATES["structured_content"] = {"html": structured_content_html, "texts": structured_content_texts}

# ── 7. big_statement ──
def big_statement_html(th, d, show):
    s = vis(show)
    stmt = (d.get("statement") or "").replace("\n","<br>")
    attr_html = f"""<div style="margin-top:30px;font-size:18px;color:{thC(th,'secondary')};visibility:{s};">— {d.get('attribution','')}</div>""" if d.get("attribution") else ""
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Georgia,serif; }}
    </style></head><body>
      <div style="position:absolute;right:-80px;top:-80px;width:500px;height:500px;border-radius:50%;background:{thC(th,'secondary')}15;"></div>
      <div style="position:absolute;left:60px;top:60px;font-size:200px;color:{thC(th,'secondary')}18;font-weight:bold;line-height:1;">"</div>
      <div style="position:absolute;left:0;top:0;right:0;bottom:0;display:flex;flex-direction:column;align-items:center;justify-content:center;padding:0 120px;text-align:center;">
        <div style="font-size:40px;font-weight:bold;color:{thC(th,'accent')};line-height:1.4;visibility:{s};">{stmt}</div>
        {attr_html}
      </div>
    </body></html>"""

def big_statement_texts(th, d):
    t = [T(120, 300, 1680, 350, (d.get("statement","") or "").replace("\n","\n"), fs=32, c=th.get("accent","FFFFFF"), ff="Georgia", b=True, a="center")]
    if d.get("attribution"):
        t.append(T(120, 680, 1680, 35, "— "+d["attribution"], fs=14, c=th.get("secondary","FFFFFF"), a="center"))
    return t

TEMPLATES["big_statement"] = {"html": big_statement_html, "texts": big_statement_texts}

# ── 8. end ──
def end_html(th, d, show):
    s = vis(show)
    title_html = (d.get("title") or "").replace("\n","<br>")
    subtitle_html = f"""<div style="margin-top:24px;font-size:24px;color:{thC(th,'secondary')};visibility:{s};">{d.get('subtitle','')}</div>""" if d.get("subtitle") else ""
    contact_html = f"""<div style="margin-top:50px;font-size:16px;color:{thC(th,'secondary')}aa;visibility:{s};">{d.get('contact','')}</div>""" if d.get("contact") else ""
    return f"""<html><head><meta charset="utf-8"><style>{RESET}
      body {{ background:{thC(th,'primary')}; font-family:Georgia,serif; }}
    </style></head><body>
      <div style="position:absolute;left:-100px;top:-100px;width:350px;height:350px;border-radius:50%;background:{thC(th,'secondary')}1a;"></div>
      <div style="position:absolute;right:-50px;bottom:-50px;width:400px;height:280px;background:{thC(th,'secondary')}22;border-radius:12px;"></div>
      <div style="position:absolute;left:0;top:0;right:0;bottom:0;display:flex;flex-direction:column;align-items:center;justify-content:center;">
        <div style="font-size:52px;font-weight:bold;color:{thC(th,'accent')};visibility:{s};">{title_html}</div>
        {subtitle_html}
        {contact_html}
      </div>
    </body></html>"""

def end_texts(th, d):
    t = [T(100, 300, 1720, 200, (d.get("title","") or "").replace("\n","\n"), fs=42, c=th.get("accent","FFFFFF"), ff="Georgia", b=True, a="center")]
    if d.get("subtitle"):
        t.append(T(100, 530, 1720, 40, d["subtitle"], fs=19, c=th.get("secondary","FFFFFF"), a="center"))
    if d.get("contact"):
        t.append(T(100, 620, 1720, 30, d["contact"], fs=13, c=th.get("secondary","FFFFFF"), a="center"))
    return t

TEMPLATES["end"] = {"html": end_html, "texts": end_texts}

# ══════════════════════════════════════════════════════════
# PPTX TEXT OVERLAY
# ══════════════════════════════════════════════════════════
ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

def add_text_overlay(slide, t_dict):
    """Add a single text box overlay to a slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR
    td = t_dict
    # Convert pixel coords → inches
    left   = Inches(px(td["x"], "x"))
    top    = Inches(px(td["y"], "y"))
    width  = Inches(px(td["w"], "x"))
    height = Inches(px(td["h"], "y"))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Vertical alignment
    va = td.get("va", "middle")
    try:
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(va, MSO_ANCHOR.MIDDLE)
    except Exception:
        pass

    text = td.get("text") or ""
    lines = str(text).split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = ALIGN_MAP.get(td.get("a","left"), PP_ALIGN.LEFT)
        run = p.runs[0] if p.runs else p.add_run()
        run.text = line
        run.font.size = Pt(round((td.get("fs") or 14) * 0.75))
        run.font.bold = bool(td.get("b"))
        run.font.italic = bool(td.get("italic"))
        # Font face — take first from comma list
        ff = str(td.get("ff") or "Calibri").split(",")[0].strip()
        run.font.name = ff
        # Color — strip # and parse
        c = str(td.get("c") or "FFFFFF").lstrip("#")[:6].ljust(6,"0")
        try:
            run.font.color.rgb = RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
        except Exception:
            pass

# ══════════════════════════════════════════════════════════
# SCREENSHOT RENDERER (Playwright)
# ══════════════════════════════════════════════════════════
def render_slide_backgrounds(slides_html, tmp_dir):
    """Render list of HTML strings to PNG files, return list of paths."""
    paths = []
    with sync_playwright() as p:
        browser = p.chromium.launch(
            executable_path="/opt/pw-browsers/chromium-1194/chrome-linux/chrome",
            args=["--no-sandbox","--disable-setuid-sandbox","--disable-dev-shm-usage","--disable-gpu"]
        )
        for i, html in enumerate(slides_html):
            page = browser.new_page(viewport={"width": W, "height": H})
            page.set_content(html, wait_until="domcontentloaded", timeout=15000)
            path = os.path.join(tmp_dir, f"bg_{i}.png")
            page.screenshot(path=path, full_page=False)
            page.close()
            paths.append(path)
        browser.close()
    return paths

# ══════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════
def main():
    args = sys.argv[1:]
    if not args:
        print("Usage: python3 html2pptx.py <input.json> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    with open(args[0], "r", encoding="utf-8") as f:
        plan = json.load(f)

    output_path = args[1] if len(args) > 1 else "output.pptx"
    th = plan.get("theme", {})
    slides_data = plan.get("slides", [])

    if not th:
        print("Error: missing 'theme' in input JSON.", file=sys.stderr)
        sys.exit(1)

    tmp_dir = tempfile.mkdtemp(prefix="html2pptx_")

    # ── Phase 1: collect background HTML for all slides ──
    print("Collecting slide HTML...")
    bg_htmls = []
    valid_slides = []
    for sd in slides_data:
        tpl = TEMPLATES.get(sd.get("template",""))
        if not tpl:
            print(f"  ⚠ Unknown template '{sd.get('template')}' (page {sd.get('page')}), skipping.")
            continue
        bg_htmls.append(tpl["html"](th, sd, False))
        valid_slides.append(sd)

    # ── Phase 2: render all backgrounds in one browser session ──
    print(f"Rendering {len(bg_htmls)} slide background(s) via Playwright...")
    bg_paths = render_slide_backgrounds(bg_htmls, tmp_dir)

    # ── Phase 3: build PPTX ──
    print("Building PPTX...")
    prs = Presentation()
    prs.slide_width  = Inches(IW)
    prs.slide_height = Inches(IH)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for i, (sd, bg_path) in enumerate(zip(valid_slides, bg_paths)):
        tpl = TEMPLATES[sd["template"]]
        print(f"  Slide {i+1}/{len(valid_slides)}: {sd['template']}")

        slide = prs.slides.add_slide(blank_layout)

        # Set background image
        with open(bg_path, "rb") as f:
            bg_data = f.read()
        bg_b64 = base64.b64encode(bg_data).decode()

        # Add background as full-slide picture
        from pptx.util import Emu
        pic = slide.shapes.add_picture(
            bg_path,
            left=Emu(0), top=Emu(0),
            width=prs.slide_width, height=prs.slide_height
        )
        # Send to back
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Add text overlays
        texts = tpl["texts"](th, sd)
        for t_dict in texts:
            if t_dict.get("text"):
                add_text_overlay(slide, t_dict)

    prs.save(output_path)

    # Cleanup
    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)

    print(f"\nDone: {output_path} ({len(valid_slides)} slides)")

if __name__ == "__main__":
    main()
