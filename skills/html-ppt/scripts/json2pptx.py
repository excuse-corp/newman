#!/usr/bin/env python3
"""
json2pptx.py — Pure-Python PPT generator (Python rewrite of json2pptx.js)

Renders all slide elements directly via python-pptx shapes + text —
no browser, no screenshots. Use as the primary generator or fallback
when Playwright is unavailable.

Usage: python3 json2pptx.py <input.json> [output.pptx]
"""

import sys
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import lxml.etree as etree

# ── Slide size ──
IW, IH = 10.0, 5.625   # inches (16:9)

def rgb(hex_str):
    """Parse hex color (with or without #) → RGBColor."""
    h = str(hex_str or "FFFFFF").lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    h = h[:6].ljust(6, "0")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def th_rgb(th, key, fallback="FFFFFF"):
    return rgb(th.get(key, fallback))

# ── Shape helpers ──
def add_rect(slide, prs, x, y, w, h, fill_hex=None, line_hex=None, line_w=1, transparency=0):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1, but we use add_shape with AUTO_SHAPE_TYPE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
        if transparency > 0:
            shape.fill.fore_color.brightness = 0  # reset
            # transparency via XML lumMod / alpha
            sp_pr = shape._element.spPr
            solid_fill = sp_pr.find(qn('a:solidFill'))
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is None:
                    srgb = etree.SubElement(solid_fill, qn('a:srgbClr'))
                    srgb.set('val', str(fill_hex).lstrip('#').ljust(6,'0').upper())
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha_val = int((1.0 - transparency/100.0) * 100000)
                alpha.set('val', str(alpha_val))
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_hex=None, transparency=0):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))  # 9 = OVAL
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
        if transparency > 0:
            sp_pr = shape._element.spPr
            solid_fill = sp_pr.find(qn('a:solidFill'))
            if solid_fill is not None:
                srgb = solid_fill.find(qn('a:srgbClr'))
                if srgb is None:
                    srgb = etree.SubElement(solid_fill, qn('a:srgbClr'))
                    srgb.set('val', str(fill_hex).lstrip('#').ljust(6,'0').upper())
                alpha = etree.SubElement(srgb, qn('a:alpha'))
                alpha_val = int((1.0 - transparency/100.0) * 100000)
                alpha.set('val', str(alpha_val))
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

def add_line(slide, x, y, w, h, color_hex, line_w=1):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(20, Inches(x), Inches(y), Inches(w) if w > 0 else Inches(0.001), Inches(h) if h > 0 else Inches(0.001))
    shape.fill.background()
    shape.line.color.rgb = rgb(color_hex)
    shape.line.width = Pt(line_w)
    return shape

ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

def add_text(slide, text, x, y, w, h, font_size=14, font_face="Calibri", color="333333",
             bold=False, italic=False, align="left", valign="top",
             word_wrap=True, transparency=0, para_space_after=0):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    va_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = va_map.get(valign, MSO_ANCHOR.TOP)

    lines = str(text or "").split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        if para_space_after:
            p.space_after = Pt(para_space_after)
        runs = p.runs
        run = runs[0] if runs else p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = str(font_face).split(",")[0].strip()
        try:
            run.font.color.rgb = rgb(color)
        except Exception:
            pass

def add_bullets(slide, items, x, y, w, h, th):
    if not items:
        return
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        # bullet via XML
        pPr = p._p.get_or_add_pPr()
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
        runs = p.runs
        run = runs[0] if runs else p.add_run()
        run.text = str(item)
        run.font.size = Pt(16)
        run.font.name = th.get("fontBody","Calibri")
        try:
            run.font.color.rgb = rgb(th.get("text_body","333333"))
        except Exception:
            pass

def add_page_badge(slide, page_num, th):
    x, y, d = 9.3, 5.1, 0.4
    add_oval(slide, x, y, d, d, fill_hex=th.get("secondary","3B7DD8"))
    add_text(slide, str(page_num), x, y, d, d,
             font_size=11, font_face=th.get("fontBody","Calibri"),
             color=th.get("accent","FFFFFF"), bold=True, align="center", valign="middle")

def accent_bar_w(th):
    style = (th.get("style") or "soft").lower()
    return 0.04 if style == "sharp" else 0.08 if style == "rounded" else 0.06

def add_title_with_accent(slide, title, th, y=0.35, bar_color=None, text_color=None):
    bw = accent_bar_w(th)
    add_rect(slide, None, 0.6, y, bw, 0.55, fill_hex=bar_color or th.get("primary","0052CC"))
    add_text(slide, title or "", 0.85, y, 8.0, 0.55,
             font_size=28, font_face=th.get("fontHeading","Georgia"),
             color=text_color or th.get("text_dark","1A1A2E"), bold=True, valign="middle")
    add_rect(slide, None, 0.6, y+0.7, 1.5, 0.03, fill_hex=th.get("secondary","3B7DD8"))

# Fix: add_rect needs prs for shapes — let's simplify by using direct shape type codes
def _shape(slide, shape_type_int, x, y, w, h):
    """Add an auto shape and return it."""
    return slide.shapes.add_shape(shape_type_int, Inches(x), Inches(y), Inches(w), Inches(h))

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1, transparency=0):
    s = _shape(slide, 1, x, y, w, h)  # 1 = RECTANGLE
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        if transparency > 0:
            _set_alpha(s, transparency)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def oval(slide, x, y, w, h, fill=None, transparency=0):
    s = _shape(slide, 9, x, y, w, h)  # 9 = OVAL
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        if transparency > 0:
            _set_alpha(s, transparency)
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def _set_alpha(shape, transparency_pct):
    """Set alpha (opacity) on a shape's solidFill via XML."""
    sp_pr = shape._element.spPr
    solid = sp_pr.find(qn('a:solidFill'))
    if solid is None:
        return
    # find or create srgbClr
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is None:
        schemeClr = solid.find(qn('a:schemeClr'))
        if schemeClr is not None:
            srgb = schemeClr
        else:
            return
    # Remove existing alpha if any
    for a in srgb.findall(qn('a:alpha')):
        srgb.remove(a)
    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
    alpha_val = int((1.0 - transparency_pct / 100.0) * 100000)
    alpha_el.set('val', str(alpha_val))

def txt(slide, text, x, y, w, h, fs=14, ff="Calibri", color="333333",
        bold=False, italic=False, align="left", valign="top", wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                           "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    lines = str(text or "").split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        run = p.runs[0] if p.runs else p.add_run()
        run.text = line
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = str(ff).split(",")[0].strip()
        try:
            run.font.color.rgb = rgb(color)
        except Exception:
            pass

def set_slide_bg(slide, color_hex):
    """Set solid slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)

# ══════════════════════════════════════════════════════════
# TEMPLATE RENDERERS
# ══════════════════════════════════════════════════════════

def render_cover(slide, d, th):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    rect(slide, 0, 0, 3.8, 5.625, fill=th.get("primary","0052CC"))
    oval(slide, -0.8, -0.8, 3.2, 3.2, fill=th.get("secondary","3B7DD8"), transparency=30)
    rect(slide, 4.4, 2.0, 1.2, 0.05, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("title",""), 4.4, 2.2, 5.0, 1.5, fs=40, ff=th.get("fontHeading","Georgia"),
        color=th.get("text_dark","1A1A2E"), bold=True, valign="top")
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 4.4, 3.7, 5.0, 0.6, fs=20, ff=th.get("fontBody","Calibri"),
            color=th.get("text_body","555555"))
    if d.get("meta"):
        txt(slide, d["meta"], 4.4, 5.0, 5.0, 0.4, fs=12, ff=th.get("fontBody","Calibri"),
            color=th.get("text_body","555555"))

def render_cover_pro(slide, d, th):
    set_slide_bg(slide, th.get("primary","0052CC"))
    # Side accent
    rect(slide, 0.4, 0.3, 0.05, 0.35, fill=th.get("secondary","3B7DD8"))
    if d.get("department"):
        txt(slide, d["department"], 0.6, 0.3, 6.0, 0.35, fs=12, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"))
    # Title
    txt(slide, d.get("title",""), 0.5, 1.3, 8.5, 2.0, fs=48, ff=th.get("fontHeading","Georgia"),
        color=th.get("accent","FFFFFF"), bold=True, valign="top")
    # Rule + subtitle
    rect(slide, 0.5, 3.5, 1.5, 0.04, fill=th.get("secondary","3B7DD8"))
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 2.2, 3.35, 6.0, 0.5, fs=22, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"))
    # Bottom line
    rect(slide, 0.4, 4.8, 9.2, 0.01, fill=th.get("secondary","3B7DD8"))
    # Stats
    for i, st in enumerate(d.get("bottom_stats") or []):
        x = 0.5 + i * 2.8
        oval(slide, x, 4.95, 0.35, 0.35, fill=th.get("secondary","3B7DD8"))
        txt(slide, st.get("text",""), x+0.45, 4.95, 2.2, 0.35, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("secondary","3B7DD8"), valign="middle")
    if d.get("date"):
        txt(slide, d["date"], 7.5, 4.95, 2.0, 0.35, fs=14, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"), align="right", valign="middle")

def render_section_break(slide, d, th):
    set_slide_bg(slide, th.get("primary","0052CC"))
    if d.get("section_number"):
        txt(slide, d["section_number"], 0.2, 0.5, 4.0, 3.5, fs=160,
            ff=th.get("fontHeading","Georgia"), color="3B7DD8",
            bold=True)
    txt(slide, d.get("title",""), 0.6, 2.4, 7.0, 0.9, fs=36, ff=th.get("fontHeading","Georgia"),
        color=th.get("accent","FFFFFF"), bold=True)
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 0.6, 3.4, 7.0, 0.5, fs=16, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"))
    rect(slide, 9.5, 0.85, 0.04, 3.94, fill=th.get("secondary","3B7DD8"))

def render_two_column(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    _add_title_accent(slide, d.get("title",""), th)
    add_page_badge(slide, pg, th)
    rect(slide, 4.65, 1.2, 5.1, 4.0, fill="FFFFFF")
    lft = d.get("left") or {}
    rgt = d.get("right") or {}
    if lft.get("type") == "bullet_list" and isinstance(lft.get("content"), list):
        add_bullets(slide, lft["content"], 0.6, 1.3, 3.8, 3.8, th)
    elif lft.get("content"):
        txt(slide, str(lft["content"]), 0.6, 1.3, 3.8, 3.8, fs=16,
            ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")
    if rgt.get("type") == "bullet_list" and isinstance(rgt.get("content"), list):
        add_bullets(slide, rgt["content"], 4.8, 1.3, 4.8, 3.8, th)
    elif rgt.get("content"):
        txt(slide, str(rgt["content"]), 4.8, 1.3, 4.8, 3.8, fs=16,
            ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")

def render_stats_callout(slide, d, th, pg):
    set_slide_bg(slide, th.get("primary","0052CC"))
    rect(slide, 0, 3.1, 10, 2.525, fill=th.get("bg_light","F5F5F5"))
    if d.get("title"):
        txt(slide, d["title"], 0.6, 0.5, 8.8, 0.6, fs=28, ff=th.get("fontHeading","Georgia"),
            color=th.get("accent","FFFFFF"), bold=True)
    add_page_badge(slide, pg, th)
    stats = d.get("stats") or []
    cnt = len(stats) or 1
    cW = 4.0 if cnt <= 2 else 2.6
    gap = 0.4
    totW = cnt * cW + (cnt-1) * gap
    sX = (10 - totW) / 2
    for i, s in enumerate(stats):
        x = sX + i*(cW+gap)
        rect(slide, x, 1.5, cW, 3.2, fill="FFFFFF")
        rect(slide, x, 1.5, cW, 0.06, fill=th.get("secondary","3B7DD8"))
        txt(slide, s.get("value",""), x, 1.8, cW, 1.2, fs=56,
            ff=th.get("fontHeading","Georgia"), color=th.get("primary","0052CC"),
            bold=True, align="center")
        txt(slide, s.get("label",""), x, 3.1, cW, 0.4, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"),
            bold=True, align="center")
        if s.get("description"):
            txt(slide, s["description"], x, 3.6, cW, 0.5, fs=12,
                ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"),
                align="center")

def render_timeline(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    rect(slide, 0, 2.1, 10, 0.7, fill=th.get("primary","0052CC"), transparency=90)
    _add_title_accent(slide, d.get("title",""), th)
    add_page_badge(slide, pg, th)
    steps = d.get("steps") or []
    cnt = len(steps) or 1
    sW = 8.8 / cnt
    for i, s in enumerate(steps):
        cx = 0.6 + sW*i + sW/2
        cY = 2.15
        if i < cnt-1:
            rect(slide, cx+0.3, cY+0.25, sW-0.6, 0.02, fill=th.get("secondary","3B7DD8"))
        oval(slide, cx-0.25, cY, 0.5, 0.5, fill=th.get("primary","0052CC"))
        txt(slide, s.get("number", str(i+1)), cx-0.25, cY, 0.5, 0.5, fs=18,
            ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"),
            bold=True, align="center", valign="middle")
        txt(slide, s.get("label",""), cx-sW/2+0.1, cY+0.65, sW-0.2, 0.4, fs=16,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"),
            bold=True, align="center")
        if s.get("description"):
            txt(slide, s["description"], cx-sW/2+0.1, cY+1.1, sW-0.2, 0.8, fs=12,
                ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"),
                align="center")

def render_comparison(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    _add_title_accent(slide, d.get("title",""), th)
    add_page_badge(slide, pg, th)
    rect(slide, 0.6, 1.3, 4.1, 3.8, fill="FFFFFF")
    rect(slide, 0.6, 1.3, 4.1, 0.06, fill=th.get("primary","0052CC"))
    txt(slide, d.get("left_label","A"), 0.85, 1.55, 3.6, 0.4, fs=20,
        ff=th.get("fontHeading","Georgia"), color=th.get("primary","0052CC"), bold=True)
    if d.get("left_items"):
        add_bullets(slide, d["left_items"], 0.85, 2.1, 3.6, 2.7, th)
    rect(slide, 5.3, 1.3, 4.1, 3.8, fill="FFFFFF")
    rect(slide, 5.3, 1.3, 4.1, 0.06, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("right_label","B"), 5.55, 1.55, 3.6, 0.4, fs=20,
        ff=th.get("fontHeading","Georgia"), color=th.get("primary","0052CC"), bold=True)
    if d.get("right_items"):
        add_bullets(slide, d["right_items"], 5.55, 2.1, 3.6, 2.7, th)

def render_big_statement(slide, d, th):
    set_slide_bg(slide, th.get("primary","0052CC"))
    oval(slide, 6, -1, 5, 5, fill=th.get("secondary","3B7DD8"), transparency=80)
    txt(slide, "\u201C", 0.8, 0.5, 2.0, 2.0, fs=160, ff=th.get("fontHeading","Georgia"),
        color=th.get("secondary","3B7DD8"), bold=True)
    txt(slide, d.get("statement",""), 1.2, 1.8, 7.6, 2.2, fs=32,
        ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"),
        bold=True, align="center", valign="middle")
    if d.get("attribution"):
        txt(slide, "— "+d["attribution"], 1.2, 4.1, 7.6, 0.5, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("secondary","3B7DD8"), align="center")

def render_image_text(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    img_side = d.get("image_side","left")
    bX = 0 if img_side == "left" else 5.5
    cX = 5.0 if img_side == "left" else 0.6
    rect(slide, bX, 0, 4.5, 5.625, fill=th.get("primary","0052CC"))
    # placeholder box for image
    rect(slide, bX+0.3, 0.5, 3.9, 4.6, fill="E8E8E8", line="BBBBBB", line_w=0.5)
    txt(slide, "建议配图：" + (d.get("image_label") or ""), bX+0.3, 0.5, 3.9, 4.6,
        fs=12, color="888888", align="center", valign="middle")
    barX = 4.5 if img_side == "left" else 5.5
    rect(slide, barX, 0.8, 0.05, 4.0, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("title",""), cX, 0.5, 4.2, 0.7, fs=28,
        ff=th.get("fontHeading","Georgia"), color=th.get("text_dark","1A1A2E"), bold=True)
    text_val = "\n\n".join(d["text"]) if isinstance(d.get("text"), list) else (d.get("text") or "")
    txt(slide, text_val, cX, 1.4, 4.2, 3.5, fs=16, ff=th.get("fontBody","Calibri"),
        color=th.get("text_body","555555"), valign="top")
    add_page_badge(slide, pg, th)

def render_end(slide, d, th):
    set_slide_bg(slide, th.get("primary","0052CC"))
    rect(slide, 6.5, 3.5, 3.5, 2.125, fill=th.get("secondary","3B7DD8"), transparency=30)
    oval(slide, -1, -1, 3, 3, fill=th.get("secondary","3B7DD8"), transparency=80)
    txt(slide, d.get("title","Thank You"), 1.0, 1.8, 8.0, 1.0, fs=40,
        ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"),
        bold=True, align="center")
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 1.0, 3.0, 8.0, 0.6, fs=20, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"), align="center")
    if d.get("contact"):
        txt(slide, d["contact"], 1.0, 4.2, 8.0, 0.4, fs=14, ff=th.get("fontBody","Calibri"),
            color=th.get("secondary","3B7DD8"), align="center")

def render_grid_content(slide, d, th, pg):
    set_slide_bg(slide, th.get("primary","0052CC"))
    rect(slide, 0.4, 0.3, 0.06, 0.5, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("title",""), 0.65, 0.3, 8.0, 0.5, fs=26,
        ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"), bold=True)
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 0.65, 0.85, 8.0, 0.3, fs=12, ff=th.get("fontBody","Calibri"),
            color=th.get("text_body","555555"))
    add_page_badge(slide, pg, th)
    lX, lW, cY, cH = 0.4, 3.8, 1.2, 3.2
    rect(slide, lX, cY, lW, cH, fill=th.get("bg_light","F5F5F5"), line=th.get("secondary","3B7DD8"), line_w=1)
    lc = d.get("left_card") or {}
    rect(slide, lX+0.15, cY+0.15, 0.55, 0.4, fill=th.get("secondary","3B7DD8"))
    txt(slide, lc.get("number","01"), lX+0.15, cY+0.15, 0.55, 0.4, fs=16,
        ff=th.get("fontHeading","Georgia"), color=th.get("primary","0052CC"),
        bold=True, align="center", valign="middle")
    txt(slide, lc.get("heading",""), lX+0.85, cY+0.15, lW-1.1, 0.4, fs=16,
        ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"),
        bold=True, valign="middle")
    txt(slide, str(lc.get("content","")), lX+0.2, cY+0.8, lW-0.4, cH-1.0, fs=13,
        ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")
    rX, rW = 4.5, 5.1
    rcs = d.get("right_cards") or []
    rcH = (cH - (len(rcs)-1)*0.2) / len(rcs) if rcs else cH
    for i, rc in enumerate(rcs):
        ry = cY + i*(rcH+0.2)
        rect(slide, rX, ry, rW, rcH, fill=th.get("bg_light","F5F5F5"),
             line=th.get("secondary","3B7DD8"), line_w=1)
        txt(slide, rc.get("heading",""), rX+0.3, ry+0.1, rW-0.5, 0.35, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"), bold=True)
        items = rc.get("items") or []
        for j, item in enumerate(items):
            iw = rW / max(len(items), 1)
            ix = rX + 0.3 + j*iw
            txt(slide, item.get("label",""), ix, ry+0.5, iw-0.3, 0.25, fs=12,
                ff=th.get("fontBody","Calibri"), color=th.get("secondary","3B7DD8"), bold=True)
            txt(slide, item.get("text",""), ix, ry+0.75, iw-0.3, rcH-1.0, fs=11,
                ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")
    if d.get("bottom_bar"):
        bb = d["bottom_bar"]
        rect(slide, 0.4, 4.65, 9.2, 0.65, fill=th.get("secondary","3B7DD8"))
        txt(slide, bb.get("heading",""), 0.65, 4.68, 3.0, 0.25, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("primary","0052CC"), bold=True)
        txt(slide, bb.get("text",""), 0.65, 4.93, 8.7, 0.3, fs=11,
            ff=th.get("fontBody","Calibri"), color=th.get("primary","0052CC"))

def render_structured_content(slide, d, th, pg):
    set_slide_bg(slide, th.get("primary","0052CC"))
    rect(slide, 0.4, 0.3, 0.06, 0.5, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("title",""), 0.65, 0.3, 8.0, 0.5, fs=26,
        ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"), bold=True)
    if d.get("subtitle"):
        txt(slide, d["subtitle"], 0.65, 0.85, 8.0, 0.3, fs=12, ff=th.get("fontBody","Calibri"),
            color=th.get("text_body","555555"))
    add_page_badge(slide, pg, th)
    cx, cy, cw, ch = 0.4, 1.2, 9.2, 2.8
    rect(slide, cx, cy, cw, ch, line=th.get("secondary","3B7DD8"), line_w=1)
    rect(slide, cx+0.15, cy+0.15, 0.55, 0.4, fill=th.get("secondary","3B7DD8"))
    txt(slide, d.get("number_badge","01"), cx+0.15, cy+0.15, 0.55, 0.4, fs=16,
        ff=th.get("fontHeading","Georgia"), color=th.get("primary","0052CC"),
        bold=True, align="center", valign="middle")
    txt(slide, d.get("heading",""), cx+0.85, cy+0.15, cw-1.2, 0.4, fs=18,
        ff=th.get("fontBody","Calibri"), color=th.get("accent","FFFFFF"),
        bold=True, valign="middle")
    secs = d.get("sections") or []
    secY0 = cy + 0.7
    secH = (ch - 1.1) / max(len(secs), 1)
    for i, sec in enumerate(secs):
        sy = secY0 + i*secH
        if i > 0:
            rect(slide, cx+0.3, sy, cw-0.6, 0.01, fill=th.get("secondary","3B7DD8"), transparency=50)
        txt(slide, sec.get("heading",""), cx+0.4, sy+0.08, cw-0.8, 0.3, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("accent","FFFFFF"), bold=True)
        txt(slide, sec.get("text",""), cx+0.4, sy+0.38, cw-0.8, secH-0.5, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")
    if d.get("callout"):
        coY = cy + ch - 0.5
        rect(slide, cx+0.3, coY, cw-0.6, 0.4, fill=th.get("bg_light","F5F5F5"),
             line=th.get("secondary","3B7DD8"), line_w=0.5)
        txt(slide, d["callout"].get("text",""), cx+0.6, coY, cw-1.2, 0.4, fs=11,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"),
            align="center", valign="middle")
    if d.get("bottom_bar"):
        bb = d["bottom_bar"]
        bbY = cy + ch + 0.2
        rect(slide, cx, bbY, cw, 0.55, fill=th.get("secondary","3B7DD8"))
        label = (bb.get("heading","") + "：" if bb.get("heading") else "") + bb.get("text","")
        txt(slide, label, cx+0.3, bbY, cw-0.6, 0.55, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("primary","0052CC"),
            valign="middle")

def render_card_grid(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    _add_title_accent(slide, d.get("title",""), th)
    add_page_badge(slide, pg, th)
    cards = d.get("cards") or []
    cols = d.get("columns") or (2 if len(cards) <= 4 else 3)
    rows = -(-len(cards) // cols)
    gX, gW, gY, gH, gapX, gapY = 0.5, 9.0, 1.3, 3.9, 0.25, 0.25
    cW = (gW - (cols-1)*gapX) / cols
    cH = (gH - (rows-1)*gapY) / rows if rows > 0 else gH
    for i, c in enumerate(cards):
        col = i % cols
        row = i // cols
        x = gX + col*(cW+gapX)
        y = gY + row*(cH+gapY)
        rect(slide, x, y, cW, cH, fill="FFFFFF", line=th.get("secondary","3B7DD8"), line_w=1)
        rect(slide, x, y, cW, 0.05, fill=th.get("primary","0052CC"))
        oval(slide, x+0.2, y+0.25, 0.4, 0.4, fill=th.get("primary","0052CC"))
        txt(slide, c.get("heading",""), x+0.75, y+0.25, cW-1.0, 0.4, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"),
            bold=True, valign="middle")
        txt(slide, c.get("text",""), x+0.2, y+0.8, cW-0.4, cH-1.1, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")

def render_kpi_dashboard(slide, d, th, pg):
    set_slide_bg(slide, th.get("bg_light","F5F5F5"))
    _add_title_accent(slide, d.get("title",""), th)
    add_page_badge(slide, pg, th)
    kpis = d.get("kpis") or []
    kC = len(kpis)
    kW = (9.0 - (kC-1)*0.2) / kC if kC else 9.0
    for i, k in enumerate(kpis):
        x = 0.5 + i*(kW+0.2)
        rect(slide, x, 1.2, kW, 1.0, fill=th.get("primary","0052CC"))
        txt(slide, k.get("value",""), x, 1.2, kW, 0.6, fs=32,
            ff=th.get("fontHeading","Georgia"), color=th.get("accent","FFFFFF"),
            bold=True, align="center")
        txt(slide, k.get("label",""), x, 1.8, kW, 0.35, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("secondary","3B7DD8"),
            align="center")
    dets = d.get("detail_cards") or []
    dC = len(dets)
    dW = (9.0 - (dC-1)*0.25) / dC if dC else 9.0
    for i, det in enumerate(dets):
        x = 0.5 + i*(dW+0.25)
        y = 2.5
        rect(slide, x, y, dW, 2.7, fill="FFFFFF", line=th.get("secondary","3B7DD8"), line_w=1)
        rect(slide, x, y, dW, 0.05, fill=th.get("secondary","3B7DD8"))
        txt(slide, det.get("heading",""), x+0.2, y+0.15, dW-0.4, 0.35, fs=14,
            ff=th.get("fontBody","Calibri"), color=th.get("text_dark","1A1A2E"), bold=True)
        txt(slide, det.get("text",""), x+0.2, y+0.55, dW-0.4, 2.0, fs=12,
            ff=th.get("fontBody","Calibri"), color=th.get("text_body","555555"), valign="top")

def _add_title_accent(slide, title, th, y=0.35):
    bw = accent_bar_w(th)
    rect(slide, 0.6, y, bw, 0.55, fill=th.get("primary","0052CC"))
    txt(slide, title or "", 0.85, y, 8.0, 0.55, fs=28,
        ff=th.get("fontHeading","Georgia"), color=th.get("text_dark","1A1A2E"),
        bold=True, valign="middle")
    rect(slide, 0.6, y+0.7, 1.5, 0.03, fill=th.get("secondary","3B7DD8"))

# ══════════════════════════════════════════════════════════
# RENDERER DISPATCH TABLE
# ══════════════════════════════════════════════════════════
RENDERERS = {
    "cover":              lambda s, d, th, pg: render_cover(s, d, th),
    "cover_pro":          lambda s, d, th, pg: render_cover_pro(s, d, th),
    "section_break":      lambda s, d, th, pg: render_section_break(s, d, th),
    "two_column":         lambda s, d, th, pg: render_two_column(s, d, th, pg),
    "stats_callout":      lambda s, d, th, pg: render_stats_callout(s, d, th, pg),
    "timeline":           lambda s, d, th, pg: render_timeline(s, d, th, pg),
    "comparison":         lambda s, d, th, pg: render_comparison(s, d, th, pg),
    "big_statement":      lambda s, d, th, pg: render_big_statement(s, d, th),
    "image_text":         lambda s, d, th, pg: render_image_text(s, d, th, pg),
    "end":                lambda s, d, th, pg: render_end(s, d, th),
    "grid_content":       lambda s, d, th, pg: render_grid_content(s, d, th, pg),
    "structured_content": lambda s, d, th, pg: render_structured_content(s, d, th, pg),
    "card_grid":          lambda s, d, th, pg: render_card_grid(s, d, th, pg),
    "kpi_dashboard":      lambda s, d, th, pg: render_kpi_dashboard(s, d, th, pg),
}

# ══════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════
def main():
    args = sys.argv[1:]
    if not args:
        print("Usage: python3 json2pptx.py <input.json> [output.pptx]", file=sys.stderr)
        sys.exit(1)

    with open(args[0], "r", encoding="utf-8") as f:
        plan = json.load(f)

    output_path = args[1] if len(args) > 1 else "output.pptx"
    theme = plan.get("theme", {})
    slides_data = plan.get("slides", [])

    if not theme:
        print("Error: missing 'theme' in input JSON.", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width  = Inches(IW)
    prs.slide_height = Inches(IH)
    blank_layout = prs.slide_layouts[6]

    for sd in slides_data:
        tmpl = sd.get("template","")
        renderer = RENDERERS.get(tmpl)
        if not renderer:
            print(f"  ⚠ Unknown template '{tmpl}' (page {sd.get('page')}), skipping.")
            continue
        slide = prs.slides.add_slide(blank_layout)
        renderer(slide, sd, theme, sd.get("page", 1))

    prs.save(output_path)
    print(f"Generated: {output_path} ({len(slides_data)} slides)")

if __name__ == "__main__":
    main()
