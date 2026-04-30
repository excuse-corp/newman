from __future__ import annotations

import html
import json
import re
import subprocess
from html.parser import HTMLParser
from pathlib import Path
from tempfile import TemporaryDirectory

from docx import Document as DocxDocument
from docx.document import Document as DocxDocumentType
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from backend.attachments.models import ParsedAttachment


MARKDOWN_BYTE_LIMIT = 10 * 1024 * 1024


class _VisibleHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._skip_depth = 0
        self._chunks: list[str] = []
        self._pending_break = False

    def handle_starttag(self, tag: str, attrs) -> None:
        normalized = tag.lower()
        if normalized in {"script", "style"}:
            self._skip_depth += 1
            return
        if normalized in {"p", "div", "section", "article", "header", "footer", "li", "tr", "br"}:
            self._pending_break = True
        if normalized in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            self._pending_break = True

    def handle_endtag(self, tag: str) -> None:
        normalized = tag.lower()
        if normalized in {"script", "style"} and self._skip_depth > 0:
            self._skip_depth -= 1
            return
        if normalized in {"p", "div", "section", "article", "header", "footer", "li", "tr"}:
            self._pending_break = True

    def handle_data(self, data: str) -> None:
        if self._skip_depth:
            return
        cleaned = re.sub(r"\s+", " ", data).strip()
        if not cleaned:
            return
        if self._pending_break and self._chunks and self._chunks[-1] != "\n":
            self._chunks.append("\n")
        self._chunks.append(cleaned)
        self._pending_break = False

    def get_text(self) -> str:
        text = "".join(self._chunks)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


def parse_attachment(path: Path) -> ParsedAttachment:
    suffix = path.suffix.lower()
    if suffix in {".doc", ".xls", ".ppt"}:
        with TemporaryDirectory(prefix="newman-attachment-convert-") as tmp:
            converted = _convert_legacy_office_document(path, Path(tmp))
            return parse_attachment(converted)

    if suffix in {".txt", ".md"}:
        return _parse_textual_attachment(path)
    if suffix == ".json":
        return _parse_json_attachment(path)
    if suffix in {".html", ".htm"}:
        return _parse_html_attachment(path)
    if suffix == ".pdf":
        return _parse_pdf_attachment(path)
    if suffix == ".docx":
        return _parse_docx_attachment(path)
    if suffix == ".xlsx":
        return _parse_xlsx_attachment(path)
    if suffix == ".pptx":
        return _parse_pptx_attachment(path)
    raise ValueError(f"暂不支持解析该文件类型: {suffix or '<none>'}")


def _parse_textual_attachment(path: Path) -> ParsedAttachment:
    content = path.read_text(encoding="utf-8", errors="replace").strip()
    paragraphs = [block.strip() for block in re.split(r"\n\s*\n", content) if block.strip()]
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    if path.suffix.lower() == ".md":
        markdown_lines.append(content or "（空内容）")
        plain_lines.append(content)
    else:
        for index, paragraph in enumerate(paragraphs or [content], start=1):
            if not paragraph:
                continue
            markdown_lines.append(f'<!-- loc: {json.dumps({"para": index, "type": "paragraph"}, ensure_ascii=False)} -->')
            markdown_lines.extend([paragraph, ""])
            plain_lines.append(paragraph)
    markdown = "\n".join(markdown_lines).strip() + "\n"
    return _finalize_parsed_attachment(markdown, "\n\n".join(plain_lines))


def _parse_json_attachment(path: Path) -> ParsedAttachment:
    raw_content = path.read_text(encoding="utf-8", errors="replace")
    try:
        payload = json.loads(raw_content)
        rendered = json.dumps(payload, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        rendered = raw_content
    markdown = f"# {path.name}\n\n```json\n{rendered.strip()}\n```\n"
    return _finalize_parsed_attachment(markdown, rendered)


def _parse_html_attachment(path: Path) -> ParsedAttachment:
    raw_content = path.read_text(encoding="utf-8", errors="replace")
    parser = _VisibleHTMLParser()
    parser.feed(raw_content)
    visible_text = html.unescape(parser.get_text())
    paragraphs = [block.strip() for block in re.split(r"\n\s*\n", visible_text) if block.strip()]
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    for index, paragraph in enumerate(paragraphs or [visible_text], start=1):
        if not paragraph:
            continue
        markdown_lines.append(f'<!-- loc: {json.dumps({"para": index, "type": "html_text"}, ensure_ascii=False)} -->')
        markdown_lines.extend([paragraph, ""])
        plain_lines.append(paragraph)
    return _finalize_parsed_attachment("\n".join(markdown_lines).strip() + "\n", "\n\n".join(plain_lines))


def _parse_pdf_attachment(path: Path) -> ParsedAttachment:
    reader = PdfReader(str(path))
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    warnings: list[str] = []
    nonempty_pages = 0
    for page_index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        nonempty_pages += 1
        markdown_lines.append(f'<!-- loc: {json.dumps({"page": page_index, "type": "page"}, ensure_ascii=False)} -->')
        markdown_lines.extend([f"## Page {page_index}", "", text, ""])
        plain_lines.append(text)
    if len(reader.pages) > 0 and nonempty_pages * 2 < len(reader.pages):
        warnings.append("PDF 可读文本页占比偏低，可能是扫描件，建议后续补做增强解析。")
    return _finalize_parsed_attachment("\n".join(markdown_lines).strip() + "\n", "\n\n".join(plain_lines), warnings=warnings)


def _parse_docx_attachment(path: Path) -> ParsedAttachment:
    document = DocxDocument(path)
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    para_index = 0
    table_index = 0
    for block in _iter_docx_blocks(document):
        if isinstance(block, DocxParagraph):
            text = block.text.strip()
            if not text:
                continue
            para_index += 1
            style_name = block.style.name if block.style is not None and block.style.name else ""
            heading_level = _docx_heading_level(style_name)
            markdown_lines.append(
                f'<!-- loc: {json.dumps({"para": para_index, "type": "heading" if heading_level else "paragraph"}, ensure_ascii=False)} -->'
            )
            if heading_level:
                markdown_lines.append(f"{'#' * heading_level} {text}")
            else:
                markdown_lines.append(text)
            markdown_lines.append("")
            plain_lines.append(text)
            continue
        if isinstance(block, DocxTable):
            table_rows = [[cell.text.strip() for cell in row.cells] for row in block.rows]
            if not any(any(cell for cell in row) for row in table_rows):
                continue
            table_index += 1
            markdown_lines.append(f'<!-- loc: {json.dumps({"table": table_index, "type": "table"}, ensure_ascii=False)} -->')
            markdown_lines.extend(_render_markdown_table(table_rows))
            markdown_lines.append("")
            plain_lines.append("\n".join(" | ".join(cell for cell in row if cell) for row in table_rows if any(cell for cell in row)))
    return _finalize_parsed_attachment("\n".join(markdown_lines).strip() + "\n", "\n\n".join(plain_lines))


def _parse_xlsx_attachment(path: Path) -> ParsedAttachment:
    workbook = load_workbook(path, data_only=True, read_only=True)
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    for worksheet in workbook.worksheets:
        if worksheet.sheet_state != "visible":
            continue
        rows: list[list[str]] = []
        for row in worksheet.iter_rows():
            values = [_stringify_excel_cell(cell.value) for cell in row]
            if any(value for value in values):
                rows.append(values)
        if not rows:
            continue
        markdown_lines.append(f'<!-- loc: {json.dumps({"sheet": worksheet.title, "type": "sheet"}, ensure_ascii=False)} -->')
        markdown_lines.extend([f"## Sheet: {worksheet.title}", ""])
        header = ["row"] + [f"col_{index}" for index in range(1, max(len(row) for row in rows) + 1)]
        table_rows = [header]
        for row_index, row_values in enumerate(rows, start=1):
            table_rows.append([str(row_index), *row_values, *[""] * (len(header) - len(row_values) - 1)])
        markdown_lines.extend(_render_markdown_table(table_rows))
        markdown_lines.append("")
        plain_lines.append(f"{worksheet.title}\n" + "\n".join(" | ".join(row) for row in table_rows[1:]))
    return _finalize_parsed_attachment("\n".join(markdown_lines).strip() + "\n", "\n\n".join(plain_lines))


def _parse_pptx_attachment(path: Path) -> ParsedAttachment:
    presentation = Presentation(path)
    markdown_lines = [f"# {path.name}", ""]
    plain_lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_title = slide.shapes.title.text.strip() if slide.shapes.title is not None and slide.shapes.title.text else f"Slide {slide_index}"
        markdown_lines.append(f'<!-- loc: {json.dumps({"slide": slide_index, "type": "slide"}, ensure_ascii=False)} -->')
        markdown_lines.extend([f"## Slide {slide_index}: {slide_title}", ""])
        plain_parts: list[str] = [slide_title]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text and text != slide_title:
                    markdown_lines.append(text)
                    markdown_lines.append("")
                    plain_parts.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                if any(any(cell for cell in row) for row in rows):
                    markdown_lines.extend(_render_markdown_table(rows))
                    markdown_lines.append("")
                    plain_parts.append("\n".join(" | ".join(cell for cell in row if cell) for row in rows if any(cell for cell in row)))
        notes_text = _read_slide_notes_text(slide)
        if notes_text:
            markdown_lines.append(f'<!-- loc: {json.dumps({"slide": slide_index, "type": "notes"}, ensure_ascii=False)} -->')
            markdown_lines.extend(["### Notes", "", notes_text, ""])
            plain_parts.append(notes_text)
        plain_lines.append("\n".join(part for part in plain_parts if part))
    return _finalize_parsed_attachment("\n".join(markdown_lines).strip() + "\n", "\n\n".join(plain_lines))


def _iter_docx_blocks(document: DocxDocumentType):
    for child in document.element.body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "p":
            yield DocxParagraph(child, document)
        elif tag == "tbl":
            yield DocxTable(child, document)


def _docx_heading_level(style_name: str) -> int | None:
    normalized = style_name.strip().lower()
    if not normalized.startswith("heading"):
        return None
    digits = "".join(char for char in normalized if char.isdigit())
    if not digits:
        return 1
    try:
        return max(1, min(int(digits), 6))
    except ValueError:
        return 1


def _render_markdown_table(rows: list[list[str]]) -> list[str]:
    width = max(len(row) for row in rows) if rows else 0
    normalized_rows = [row + [""] * (width - len(row)) for row in rows]
    if not normalized_rows:
        return []
    header = normalized_rows[0]
    separator = ["---"] * width
    lines = [
        "| " + " | ".join(cell or " " for cell in header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in normalized_rows[1:]:
        lines.append("| " + " | ".join(cell or " " for cell in row) + " |")
    return lines


def _stringify_excel_cell(value: object) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _read_slide_notes_text(slide) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        text = slide.notes_slide.notes_text_frame.text or ""
        return text.strip()
    except Exception:
        return ""


def _convert_legacy_office_document(path: Path, output_dir: Path) -> Path:
    target_extension = {
        ".doc": "docx",
        ".xls": "xlsx",
        ".ppt": "pptx",
    }[path.suffix.lower()]
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        target_extension,
        "--outdir",
        str(output_dir),
        str(path),
    ]
    try:
        subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=True,
            text=True,
            timeout=120,
        )
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
        raise ValueError(f"{path.name} 转换失败，无法解析旧版 Office 文件") from exc
    converted = output_dir / f"{path.stem}.{target_extension}"
    if not converted.exists():
        matches = list(output_dir.glob(f"*.{target_extension}"))
        if not matches:
            raise ValueError(f"{path.name} 转换失败，无法找到转换结果")
        converted = matches[0]
    return converted.resolve()


def _finalize_parsed_attachment(markdown: str, plain_text: str, *, warnings: list[str] | None = None) -> ParsedAttachment:
    normalized_markdown = markdown.strip() + "\n"
    if len(normalized_markdown.encode("utf-8")) > MARKDOWN_BYTE_LIMIT:
        raise ValueError("解析结果过大，当前附件不适合直接注入对话上下文")
    meaningful_text = re.sub(r"\s+", " ", plain_text).strip()
    if len(meaningful_text) < 8:
        raise ValueError("正文抽取不足，当前附件无法形成稳定解析结果")
    return ParsedAttachment(markdown=normalized_markdown, plain_text=meaningful_text, warnings=warnings or [])
